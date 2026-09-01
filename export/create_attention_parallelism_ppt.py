#!/usr/bin/env python3
"""Create a one-slide TP8 MLA attention parallelism overview."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from create_moe_route_ppt import (
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5


def add_flow_box(
    slide,
    x,
    y,
    w,
    h,
    title,
    formula,
    detail,
    accent,
    title_size=8.8,
    formula_size=7.5,
    detail_size=7.1,
):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "panel", "line", 0.9)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.065, h, accent)
    add_text(
        slide,
        title,
        x + 0.13,
        y + 0.10,
        w - 0.22,
        0.25,
        title_size,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        formula,
        x + 0.10,
        y + 0.43,
        w - 0.18,
        0.28,
        formula_size,
        accent,
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        x + 0.12,
        y + 0.78,
        w - 0.22,
        h - 0.84,
        detail_size,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.TOP,
    )


def add_right_arrow(slide, x, y, color="line"):
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x, y, 0.14, 0.18, color)


def add_down_arrow(slide, x, y, color="line"):
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, x, y, 0.18, 0.14, color)


def add_info_card(slide, x, y, w, h, title, body, accent, body_size=7.4):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "light", accent, 0.9)
    add_text(
        slide,
        title,
        x + 0.14,
        y + 0.10,
        w - 0.28,
        0.24,
        9.0,
        accent,
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        body,
        x + 0.14,
        y + 0.40,
        w - 0.28,
        h - 0.49,
        body_size,
        "black",
        False,
        CODE_FONT if "[" in body or "P_" in body else FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("F6F9FB")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.045, "cyan")
    add_text(
        slide,
        "MLA Attention：TP8 并行路径与计算公式",
        0.35,
        0.13,
        9.35,
        0.40,
        21.0,
        "red",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Q heads / O factors 按 TP 切分；latent KV 每 rank 复制；8 个 head 的结果先本地聚合，再对 wo_b 部分结果做 AllReduce",
        0.38,
        0.58,
        10.15,
        0.22,
        8.9,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.76, 0.18, 2.18, 0.38, "panel_2", "blue", 0.8)
    add_text(
        slide,
        "H=64  d=512  G=8  TP=8",
        10.85,
        0.27,
        2.00,
        0.20,
        8.1,
        "blue",
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.35, 0.91, 12.60, 0.35, "panel_2", "line", 0.8)
    add_text(slide, "Replicate", 0.54, 0.995, 0.77, 0.18, 8.0, "cyan", True, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "每 rank 完整副本", 1.31, 0.995, 1.32, 0.18, 7.8, "black", False, FONT_CJK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "Shard", 2.89, 0.995, 0.54, 0.18, 8.0, "blue", True, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "按 head / group / input-output dim 切分", 3.43, 0.995, 2.68, 0.18, 7.8, "black", False, FONT_CJK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "Local", 6.35, 0.995, 0.54, 0.18, 8.0, "green", True, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "本 rank 计算，不发生逐 head 网络传输", 6.89, 0.995, 2.60, 0.18, 7.8, "black", False, FONT_CJK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "AllReduce", 9.85, 0.995, 0.92, 0.18, 8.0, "purple", True, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "各 rank partial output 求和", 10.77, 0.995, 1.86, 0.18, 7.8, "black", False, FONT_CJK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

    add_text(slide, "A  Q path + local Sparse Attention", 0.38, 1.36, 4.20, 0.20, 8.0, "gray", True, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    flow_y, flow_h = 1.62, 1.30
    boxes = [
        (0.35, 1.28, "Input x", "[B,S,D]", "D=4096\nreplicated", "cyan"),
        (1.75, 1.53, "wq_a + q_norm", "D → q=1024", "full q_rank on\nevery rank", "blue"),
        (3.42, 1.45, "wq_b", "q → 8×512", "ColumnParallel\n8 local Q heads", "blue"),
        (5.00, 1.91, "Sparse Attention", "Q_local × KV", "QK + AV\n8 heads, same KV candidates", "green"),
        (7.04, 1.55, "reshape / group", "[B,S,8,512]", "→ [B,S,1,4096]\n8 heads → 1 group", "amber"),
        (8.76, 1.50, "wo_a", "4096 → 1024", "local output factor\nper group", "amber"),
        (10.43, 1.42, "wo_b", "1024 → D", "partial [B,S,D]\nRowParallel", "purple"),
        (12.00, 0.94, "AllReduce", "Σ rank", "final x\n[B,S,D]", "purple"),
    ]
    for x, w, title, formula, detail, accent in boxes:
        add_flow_box(slide, x, flow_y, w, flow_h, title, formula, detail, accent)
    arrow_positions = [1.66, 3.33, 4.91, 6.95, 8.67, 10.34, 11.91]
    for x in arrow_positions:
        add_right_arrow(slide, x, 2.18)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.35, 2.99, 12.60, 0.91, "panel", "line", 0.9)
    add_text(slide, "KV path", 0.53, 3.10, 0.76, 0.20, 8.8, "cyan", True, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "wkv + Compressor", 1.32, 3.08, 1.52, 0.22, 8.0, "black", True, FONT_CJK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "[B,S,D] → [B,S,512]", 1.32, 3.35, 1.70, 0.20, 7.4, "cyan", True, CODE_FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "raw W=128 + compressed KV\nreplicated on every rank", 3.10, 3.14, 2.10, 0.45, 7.2, "muted", False, FONT_CJK, PP_ALIGN.LEFT, MSO_ANCHOR.TOP)
    add_down_arrow(slide, 5.84, 2.89, "cyan")
    add_text(slide, "candidate KV", 5.58, 3.70, 0.78, 0.17, 6.8, "cyan", True, CODE_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 5.46, 3.10, 0.012, 0.60, "line")
    add_text(slide, "ratio-4 Indexer (short layers only)", 6.02, 3.08, 2.47, 0.22, 8.0, "blue", True, FONT_CJK, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "64 heads × 128 dim", 6.02, 3.35, 1.54, 0.18, 7.3, "blue", True, CODE_FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "scan P_index → score AllReduce → Top-K=512", 7.68, 3.31, 3.17, 0.24, 7.2, "muted", False, CODE_FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(slide, "topk_idxs", 10.98, 3.35, 0.83, 0.18, 7.2, "blue", True, CODE_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_down_arrow(slide, 11.29, 2.89, "blue")

    add_text(
        slide,
        "关键：多个 Q head 的结果不是跨 rank 逐个发送；它们在本 rank 内按 group 拼接后进入 wo_a，只有 wo_b 的 partial output 需要 AllReduce。",
        0.42,
        4.05,
        12.42,
        0.20,
        8.1,
        "red",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )

    card_y, card_h = 4.36, 1.15
    add_info_card(
        slide,
        0.35,
        card_y,
        4.00,
        card_h,
        "TP8 shape ledger",
        "x         [B,S,4096]       replicate\nq_rank    [B,S,1024]       replicate\nQ_local   [B,S,8,512]      shard\nKV bank   [B,P,512]        replicate",
        "cyan",
        7.2,
    )
    add_info_card(
        slide,
        4.55,
        card_y,
        4.00,
        card_h,
        "Head results → O path",
        "o_local   [B,S,8,512]\nreshape   [B,S,1,4096]   (H/G=8)\nwo_a      [B,S,1,1024]\nwo_b      [B,S,1024] → partial [B,S,D]",
        "amber",
        7.2,
    )
    add_info_card(
        slide,
        8.75,
        card_y,
        4.20,
        card_h,
        "Candidate paths",
        "P_raw   = causal raw window\nP_short = ratio-4, capped by Top-K\nP_long  = ratio-128 compressed KV\nP_index = uncapped ratio-4 scan",
        "blue",
        7.2,
    )

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.35, 5.76, 12.60, 1.00, "panel_2", "line", 0.9)
    add_text(slide, "B  FLOPs model", 0.54, 5.88, 1.28, 0.20, 8.4, "black", True, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    add_text(
        slide,
        "P_total = Lw·P_raw + Ls·(P_raw+P_short) + Ll·(P_raw+P_long)\n"
        "F_sparse,global = 4·B·H·d·P_total    |    F_sparse,rank = 4·B·max(H/TP,16)·d·P_total\n"
        "F_proj,rank = 2·B·S·L·[Dq + q(H/TP)d + Dd + (H/TP)dr + (G/TP)rD]",
        1.92,
        5.83,
        10.74,
        0.76,
        7.4,
        "blue",
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        "W=128 · H=64 · d=512 · G=8 · TP=8 · kernel_min_heads=16",
        0.54,
        6.84,
        7.20,
        0.18,
        7.0,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "source  inference/model.py  ·  calculate/generate_calculator.py",
        8.04,
        6.84,
        4.91,
        0.18,
        7.0,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "MLA Attention TP8 Parallelism"
    prs.core_properties.subject = "DeepSeek V4 Flash attention tensor parallelism and FLOPs formulas"
    prs.core_properties.author = "DeepSeek V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide overview generated from inference/model.py and calculate/generate_calculator.py."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("attention_parallelism_tp8.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()