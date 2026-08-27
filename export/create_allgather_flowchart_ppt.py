#!/usr/bin/env python3
"""Create a simple one-slide remote-write AllGather flowchart."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from create_moe_route_ppt import (
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5


def add_flow_node(
    slide,
    shape_type,
    x,
    y,
    w,
    h,
    text,
    size=12.0,
    font=FONT_CJK,
    fill="panel",
):
    add_shape(slide, shape_type, x, y, w, h, fill, "black", 1.35)
    add_text(
        slide,
        text,
        x + 0.10,
        y + 0.06,
        w - 0.20,
        h - 0.12,
        size,
        "black",
        False,
        font,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_down_arrow(slide, y):
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, 4.08, y, 0.25, 0.15, "black")


def add_legend_row(slide, symbol, meaning, y, symbol_size=10.0):
    add_text(
        slide,
        symbol,
        8.80,
        y,
        1.16,
        0.28,
        symbol_size,
        "black",
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        meaning,
        9.92,
        y,
        2.55,
        0.28,
        8.4,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_text(
        slide,
        "8 核远程写入 AllGather 流程图",
        0.48,
        0.19,
        12.25,
        0.38,
        20.0,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Core_i 视角：8 个核心同时执行同一个 AllGather 算子",
        0.51,
        0.61,
        7.40,
        0.20,
        9.0,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    main_x, main_w = 0.62, 7.58
    main_center = main_x + main_w / 2
    del main_center

    add_flow_node(
        slide,
        MSO_SHAPE.OVAL,
        2.08,
        0.95,
        5.24,
        0.56,
        "开始\n8 个核心同时启动 AllGather",
        11.5,
    )
    add_down_arrow(slide, 1.56)
    add_flow_node(
        slide,
        MSO_SHAPE.RECTANGLE,
        1.82,
        1.75,
        5.76,
        0.52,
        "Core_i 获取本地分片：  X_i  [C bytes]",
        11.3,
        CODE_FONT,
    )
    add_down_arrow(slide, 2.31)
    add_flow_node(
        slide,
        MSO_SHAPE.RECTANGLE,
        1.50,
        2.50,
        6.37,
        0.48,
        "准备接收缓冲区：  Y_i[0:8]    （8 个 slot，每个 C bytes）",
        10.6,
        CODE_FONT,
    )
    add_down_arrow(slide, 3.04)
    add_flow_node(
        slide,
        MSO_SHAPE.RECTANGLE,
        2.08,
        3.23,
        5.24,
        0.53,
        "本地放置自己的分片：  Y_i[i]  ←  X_i",
        11.0,
        CODE_FONT,
    )
    add_down_arrow(slide, 3.77)
    add_flow_node(
        slide,
        MSO_SHAPE.RECTANGLE,
        0.94,
        3.96,
        7.52,
        0.63,
        "向其他 7 个核心远程写入：\nY_j[i]  ←  X_i，  j ∈ {0,…,7} 且 j ≠ i",
        10.8,
        CODE_FONT,
    )
    add_down_arrow(slide, 4.65)
    add_flow_node(
        slide,
        MSO_SHAPE.RECTANGLE,
        2.18,
        4.83,
        5.04,
        0.52,
        "fence  core_i",
        12.0,
        CODE_FONT,
    )
    add_down_arrow(slide, 5.40)
    add_flow_node(
        slide,
        MSO_SHAPE.RECTANGLE,
        1.66,
        5.58,
        6.08,
        0.56,
        "sync  (core_0, …, core_7)",
        11.8,
        CODE_FONT,
    )
    add_down_arrow(slide, 6.22)
    add_flow_node(
        slide,
        MSO_SHAPE.RECTANGLE,
        1.94,
        6.40,
        5.52,
        0.53,
        "读取完整接收缓冲区：  Y_i[0:8] = [X_0, X_1, …, X_7]",
        9.6,
        CODE_FONT,
    )
    add_down_arrow(slide, 6.90)
    add_flow_node(
        slide,
        MSO_SHAPE.OVAL,
        2.28,
        7.08,
        4.84,
        0.36,
        "结束：AllGather 完成",
        10.8,
    )

    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.66, 0.95, 4.04, 5.96, "light", "black", 1.0)
    add_text(
        slide,
        "符号说明",
        8.84,
        1.16,
        3.68,
        0.30,
        13.5,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "8 个核心编号：core_0 … core_7",
        8.84,
        1.55,
        3.48,
        0.25,
        8.7,
        "black",
        False,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_legend_row(slide, "i", "当前核心编号，i=0,…,7", 1.98)
    add_legend_row(slide, "j", "目标核心编号，j≠i", 2.39)
    add_legend_row(slide, "C", "一个分片的字节数", 2.80)
    add_legend_row(slide, "X_i", "Core_i 持有的本地分片", 3.21)
    add_legend_row(slide, "Y_i[k]", "Core_i 缓冲区第 k 个 slot", 3.62, 9.2)
    add_legend_row(slide, "Y_j[i]", "Core_j 的第 i 个接收 slot", 4.03, 9.2)
    add_legend_row(slide, "←", "数据写入方向", 4.44, 12.0)
    add_legend_row(slide, "fence", "远程写入有序并对后续可见", 4.85, 9.0)
    add_legend_row(slide, "sync", "8 个核心的集合屏障", 5.26, 9.0)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.84, 5.78, 3.60, 0.012, "black")
    add_text(
        slide,
        "AllGather：只搬运并拼接分片，\n不做数值相加。",
        8.84,
        5.97,
        3.48,
        0.54,
        9.2,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        "单个 Core_i：本地 1 次放置 + 远程 7 次写入",
        8.84,
        6.58,
        3.48,
        0.24,
        8.2,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "8 核远程写入 AllGather 流程图"
    prs.core_properties.subject = "Simple top-down AllGather flowchart with symbol definitions"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable black-and-white one-slide flowchart for the remote-write AllGather path."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("allgather_remote_write_flowchart.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()