#!/usr/bin/env python3
"""Create a two-slide editable Transformer Block statistics summary."""

from __future__ import annotations

import sys
from dataclasses import dataclass, replace
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "calculate"))

from generate_calculator import (  # noqa: E402
    Inputs,
    LayerCounts,
    cache_values,
    count_layers,
    load_inputs,
    parameter_components,
    scenario_items,
)
from create_moe_route_ppt import (  # noqa: E402
    CODE_FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5
TP_VALUES = (1, 8)

MONO = {
    "background": "F5F9FA",
    "panel": "FFFFFF",
    "band": "E3EFF2",
    "band_alt": "EDF5F6",
    "header": "365E69",
    "accent": "2C7180",
    "accent_dark": "1E5360",
    "ink": "18323A",
    "muted": "55717A",
    "line": "A9C0C7",
    "white": "FFFFFF",
}


@dataclass(frozen=True)
class SectionSpec:
    key: str
    title: str
    subtitle: str
    config: Inputs
    layers: LayerCounts
    note: str


@dataclass(frozen=True)
class RowData:
    tp: int
    params: str
    cache: str
    dtype: str
    compute: str
    note: str


def format_count(value: float) -> str:
    if abs(value) >= 1e9:
        return f"{value / 1e9:.3f}B"
    if abs(value) >= 1e6:
        return f"{value / 1e6:.3f}M"
    if abs(value) >= 1e3:
        return f"{value / 1e3:.3f}K"
    return f"{value:.0f}"


def format_gb(value: float) -> str:
    return f"{value / 1e9:.3f} GB"


def format_flops(value: float, mode: str) -> str:
    if mode == "prefill":
        return f"{value / 1e12:.3f} TFLOPs"
    return f"{value / 1e9:.3f} GFLOPs"


def block_parameter_values(config: Inputs, layers: LayerCounts) -> tuple[float, float]:
    block_other_names = {"Hyper-Connections", "Norms and attention sinks"}
    components = parameter_components(config, layers)
    selected = [
        component
        for component in components
        if component.category in {"Attention", "MoE"}
        or component.name in block_other_names
    ]
    return (
        sum(component.rank_count for component in selected),
        sum(component.rank_bytes for component in selected),
    )


def block_compute_values(
    config: Inputs, layers: LayerCounts, mode: str
) -> tuple[float, float, float, float]:
    items, _ = scenario_items(config, layers, mode)
    attention = sum(item.rank_flops for item in items if item.category == "Attention")
    moe = sum(item.rank_flops for item in items if item.category == "MoE")
    other = sum(
        item.rank_flops
        for item in items
        if item.category == "Other" and item.name != "LM Head"
    )
    return attention, moe, other, attention + moe + other


def cache_text(config: Inputs, layers: LayerCounts, context: int, batch: int) -> str:
    cache = cache_values(config, layers, context, batch)
    return (
        f"合计 {format_gb(cache['total'])}\n"
        f"主 KV {format_gb(cache['main'])} · Indexer {format_gb(cache['indexer'])}\n"
        f"State {format_gb(cache['states'])}"
    )


def parameter_text(config: Inputs, layers: LayerCounts) -> str:
    count, storage = block_parameter_values(config, layers)
    return f"logical {format_count(count)}\n容量 {format_gb(storage)}"


def dtype_text(section_key: str) -> str:
    if section_key == "ratio4":
        return (
            "Attn FP8 + scale；wo_a BF16\n"
            "Indexer FP8 / BF16 / FP32\n"
            "MoE FP4/FP8；KV/Comp/HC/Norm FP32"
        )
    if section_key == "ratio8":
        return (
            "Attn FP8 + scale；wo_a BF16\n"
            "无 Indexer；长压缩 ratio=8\n"
            "MoE FP4/FP8；KV/Comp/HC/Norm FP32"
        )
    if section_key == "window":
        return (
            "Attn FP8 + scale；wo_a BF16\n"
            "raw KV / activation BF16\n"
            "MoE FP4/FP8；Norm/HC/Router FP32"
        )
    return (
        "Attn FP8 + scale；wo_a BF16\n"
        "MoE routed FP4 / shared FP8\n"
        "KV/activation BF16；Comp/Norm/HC/Router FP32"
    )


def compute_text(config: Inputs, layers: LayerCounts, mode: str) -> str:
    attention, moe, other, total = block_compute_values(config, layers, mode)
    return (
        f"总计 {format_flops(total, mode)}\n"
        f"Attn {format_flops(attention, mode)} · MoE {format_flops(moe, mode)}\n"
        f"HC/Norm {format_flops(other, mode)}"
    )


def make_sections(config_pair: tuple[Inputs, list[int]]) -> list[SectionSpec]:
    base, ratios = config_pair
    native_layers = count_layers(ratios, base)
    return [
        SectionSpec(
            "full",
            "整网 Transformer Block",
            "原生 43 层：window×2 + ratio=4×21 + ratio=128×20",
            base,
            native_layers,
            "43 layers\n原生 mix",
        ),
        SectionSpec(
            "ratio4",
            "ratio=4 典型层",
            "短压缩层：含 ratio-4 Indexer 与 Top-K=512",
            base,
            LayerCounts(total=1, window=0, short=1, long=0, hash=0),
            "1 layer\nshort + Indexer",
        ),
        SectionSpec(
            "ratio8",
            "ratio=8 典型层",
            "长压缩公式对比层；原生模型未启用 ratio=8",
            replace(base, long_ratio=8),
            LayerCounts(total=1, window=0, short=0, long=1, hash=0),
            "1 layer\ncomparison",
        ),
        SectionSpec(
            "window",
            "其他典型层（window）",
            "原生滑动窗口层：window=128，无压缩 KV",
            base,
            LayerCounts(total=1, window=1, short=0, long=0, hash=0),
            "1 layer\nraw window",
        ),
    ]


def add_table_cell(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    fill: str = MONO["panel"],
    color: str = MONO["ink"],
    size: float = 6.4,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = FONT_CJK,
) -> None:
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, MONO["line"], 0.55)
    add_text(
        slide,
        text,
        x + 0.07,
        y + 0.035,
        w - 0.14,
        h - 0.07,
        size,
        color,
        bold,
        font,
        align,
        valign,
    )


def add_section_band(
    slide: Any, x: float, y: float, w: float, h: float, section: SectionSpec
) -> None:
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, MONO["band"], MONO["line"], 0.6)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, h, MONO["accent"])
    add_text(
        slide,
        section.title,
        x + 0.15,
        y + 0.025,
        2.35,
        h - 0.05,
        8.4,
        MONO["accent_dark"],
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        section.subtitle,
        x + 2.55,
        y + 0.03,
        w - 2.7,
        h - 0.06,
        6.35,
        MONO["muted"],
        False,
        FONT_CJK,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )


def build_rows(
    section: SectionSpec, mode: str, base: Inputs
) -> list[RowData]:
    context = base.prefill_sequence if mode == "prefill" else base.decode_context
    batch = base.prefill_batch if mode == "prefill" else base.decode_batch
    rows = []
    for tp in TP_VALUES:
        config = replace(section.config, tp=tp)
        local_heads = int(config.heads / tp)
        local_groups = int(config.o_groups / tp)
        rows.append(
            RowData(
                tp=tp,
                params=parameter_text(config, section.layers),
                cache=cache_text(config, section.layers, context, batch),
                dtype=dtype_text(section.key),
                compute=compute_text(config, section.layers, mode),
                note=(
                    f"{section.note}\n"
                    f"Q {local_heads} heads · O {local_groups} groups"
                ),
            )
        )
    return rows


def add_summary_table(slide: Any, mode: str, config_pair: tuple[Inputs, list[int]]) -> None:
    table_x, table_y = 0.25, 0.93
    column_widths = [0.90, 2.35, 2.20, 3.25, 2.85, 1.25]
    table_w = sum(column_widths)
    header_h, section_h, row_h = 0.38, 0.28, 0.58
    sections = make_sections(config_pair)
    table_h = header_h + len(sections) * (section_h + 2 * row_h)

    add_shape(slide, MSO_SHAPE.RECTANGLE, table_x, table_y, table_w, table_h, MONO["panel"], MONO["line"], 0.8)
    headers = [
        "TP / 每卡",
        "Transformer Block\n算子参数量",
        "每卡 KV Cache",
        "dtype（参数 / 激活）",
        "计算量 / 卡",
        "层型 / 口径",
    ]
    cursor_x = table_x
    for width, header in zip(column_widths, headers):
        add_table_cell(
            slide,
            cursor_x,
            table_y,
            width,
            header_h,
            header,
            fill=MONO["header"],
            color=MONO["white"],
            size=7.0,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        cursor_x += width

    y = table_y + header_h
    base = config_pair[0]
    for section in sections:
        add_section_band(slide, table_x, y, table_w, section_h, section)
        y += section_h
        for row in build_rows(section, mode, base):
            cursor_x = table_x
            add_table_cell(
                slide,
                cursor_x,
                y,
                column_widths[0],
                row_h,
                f"TP{row.tp}\n每卡",
                fill=MONO["band_alt"],
                color=MONO["accent_dark"],
                size=8.2,
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                font=CODE_FONT,
            )
            cursor_x += column_widths[0]
            add_table_cell(slide, cursor_x, y, column_widths[1], row_h, row.params, size=6.55)
            cursor_x += column_widths[1]
            add_table_cell(slide, cursor_x, y, column_widths[2], row_h, row.cache, size=5.95)
            cursor_x += column_widths[2]
            add_table_cell(slide, cursor_x, y, column_widths[3], row_h, row.dtype, size=5.65)
            cursor_x += column_widths[3]
            add_table_cell(slide, cursor_x, y, column_widths[4], row_h, row.compute, size=5.95)
            cursor_x += column_widths[4]
            add_table_cell(
                slide,
                cursor_x,
                y,
                column_widths[5],
                row_h,
                row.note,
                fill=MONO["band_alt"],
                color=MONO["muted"],
                size=5.45,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
            )
            y += row_h


def build_slide(
    prs: Presentation, mode: str, config_pair: tuple[Inputs, list[int]]
) -> None:
    is_prefill = mode == "prefill"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(MONO["background"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.045, MONO["accent"])
    add_text(
        slide,
        f"Transformer Block 统计 · {'Prefill' if is_prefill else 'Decode'}",
        0.30,
        0.12,
        8.9,
        0.34,
        19.0,
        MONO["accent_dark"],
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    scenario_label = "B=1 · S=8,192" if is_prefill else "B=1 · 1 token · C=1,048,576"
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.38, 0.13, 2.57, 0.32, MONO["band"], MONO["accent"], 0.7)
    add_text(
        slide,
        scenario_label,
        10.48,
        0.20,
        2.37,
        0.17,
        7.0,
        MONO["accent_dark"],
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "统一口径：参数量、每卡 KV Cache、dtype、Transformer Block 计算量；TP1/TP8 均为单 Rank",
        0.33,
        0.55,
        9.65,
        0.19,
        7.65,
        MONO["muted"],
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "D=4096 · H=64 · d=512 · q=1024 · G=8",
        10.03,
        0.55,
        2.92,
        0.19,
        7.0,
        MONO["muted"],
        True,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    add_summary_table(slide, mode, config_pair)
    add_text(
        slide,
        "Block = Attention + MoE + HC/Norm；不含 Embedding、LM Head、尾部 HC。ratio=8 是公式对比假设；其他典型层取原生 window=128，原生 ratio=128 已计入整网。",
        0.32,
        7.16,
        12.55,
        0.18,
        5.9,
        MONO["muted"],
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def build_deck(output: Path) -> None:
    config_pair = load_inputs(ROOT / "inference" / "config.json")
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    build_slide(prs, "prefill", config_pair)
    build_slide(prs, "decode", config_pair)
    prs.core_properties.title = "DeepSeek V4 Flash Transformer Block Prefill Decode Summary"
    prs.core_properties.subject = "TP1/TP8 Transformer Block parameters, KV cache, dtype, and compute"
    prs.core_properties.author = "DeepSeek V4 Flash repository"
    prs.core_properties.comments = "Editable two-slide summary generated from the repository calculator."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("transformer_block_prefill_decode_summary.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()