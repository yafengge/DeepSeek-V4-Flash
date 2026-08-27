#!/usr/bin/env python3
"""Create two simple eight-core reduction flowcharts in one presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from create_moe_route_ppt import (
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5
MAIN_X, MAIN_W = 0.62, 7.58
LEGEND_X, LEGEND_Y, LEGEND_W, LEGEND_H = 8.66, 0.92, 4.04, 6.37


def add_flow_box(
    slide,
    x,
    y,
    w,
    h,
    title,
    formula,
    title_size=11.4,
    formula_size=9.0,
):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel", "black", 1.25)
    add_text(
        slide,
        title,
        x + 0.12,
        y + 0.04,
        w - 0.24,
        0.23,
        title_size,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        formula,
        x + 0.12,
        y + 0.28,
        w - 0.24,
        h - 0.32,
        formula_size,
        "black",
        False,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_endpoint(slide, x, y, w, h, title, detail, title_size=12.0):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, w, h, "panel", "black", 1.35)
    add_text(
        slide,
        title,
        x + 0.12,
        y + 0.07,
        w - 0.24,
        0.22,
        title_size,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        x + 0.14,
        y + 0.28,
        w - 0.28,
        h - 0.34,
        8.4,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_down_arrow(slide, y):
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, 4.08, y, 0.25, 0.15, "black")


def add_legend_row(slide, symbol, meaning, y, symbol_size=9.8):
    add_text(
        slide,
        symbol,
        LEGEND_X + 0.20,
        y,
        1.08,
        0.25,
        symbol_size,
        "black",
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        meaning,
        LEGEND_X + 1.29,
        y,
        LEGEND_W - 1.52,
        0.25,
        8.0,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def add_legend_panel(slide, rows, notes):
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        LEGEND_X,
        LEGEND_Y,
        LEGEND_W,
        LEGEND_H,
        "light",
        "black",
        1.0,
    )
    add_text(
        slide,
        "符号说明",
        LEGEND_X + 0.20,
        LEGEND_Y + 0.20,
        LEGEND_W - 0.40,
        0.30,
        13.0,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "8 个核心：core_0 … core_7",
        LEGEND_X + 0.20,
        LEGEND_Y + 0.60,
        LEGEND_W - 0.40,
        0.24,
        8.5,
        "black",
        False,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    row_y = LEGEND_Y + 1.02
    for row in rows:
        symbol, meaning, *size = row
        add_legend_row(slide, symbol, meaning, row_y, size[0] if size else 9.8)
        row_y += 0.38
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        LEGEND_X + 0.20,
        row_y + 0.03,
        LEGEND_W - 0.40,
        0.012,
        "black",
    )
    note_y = row_y + 0.22
    for note in notes:
        add_text(
            slide,
            note,
            LEGEND_X + 0.20,
            note_y,
            LEGEND_W - 0.40,
            0.44,
            8.2,
            "black",
            False,
            FONT_CJK,
            PP_ALIGN.LEFT,
            MSO_ANCHOR.TOP,
        )
        note_y += 0.55


def add_header(slide, title, subtitle, page_number):
    add_text(
        slide,
        title,
        0.48,
        0.18,
        10.90,
        0.39,
        20.0,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        subtitle,
        0.51,
        0.61,
        11.35,
        0.20,
        8.9,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        page_number,
        11.98,
        0.25,
        0.76,
        0.20,
        8.0,
        "gray",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )


def add_common_footer(slide, text):
    add_text(
        slide,
        text,
        0.62,
        7.30,
        12.08,
        0.14,
        6.8,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")
    return slide


def build_allreduce_slide(prs):
    slide = new_slide(prs)
    add_header(
        slide,
        "8 核 AllReduce（不分片）流程图",
        "每个 core 持有完整 tensor；先通过远程写收齐 8 份，再在目标 core 本地逐元素求和",
        "1 / 2",
    )
    add_endpoint(
        slide,
        1.56,
        0.95,
        5.72,
        0.53,
        "开始",
        "Core_0 … Core_7 同时启动 AllReduce",
        11.6,
    )
    for y in (1.55, 2.28, 3.01, 3.81, 4.52, 5.23, 5.94, 6.66):
        add_down_arrow(slide, y)
    add_flow_box(
        slide,
        1.28,
        1.74,
        6.22,
        0.46,
        "每个 Core 准备完整输入",
        "X_i  [S]",
        formula_size=9.8,
    )
    add_flow_box(
        slide,
        0.98,
        2.47,
        6.82,
        0.49,
        "为 8 份输入预留本地接收槽位",
        "A_i[0:8]  ·  每槽 S bytes",
        formula_size=8.8,
    )
    add_flow_box(
        slide,
        1.50,
        3.20,
        5.78,
        0.49,
        "放置本地副本",
        "A_i[i]  ←  X_i",
        formula_size=9.8,
    )
    add_flow_box(
        slide,
        0.72,
        4.00,
        7.34,
        0.61,
        "远程写入其他 7 个核心",
        "remote_write(X_i  →  A_j[i], S)   ·   j ≠ i",
        title_size=11.0,
        formula_size=8.5,
    )
    add_flow_box(
        slide,
        1.68,
        4.71,
        5.42,
        0.47,
        "本核写入顺序栅栏",
        "fence core_i",
        formula_size=10.0,
    )
    add_flow_box(
        slide,
        1.42,
        5.42,
        5.94,
        0.51,
        "8 核集合同步",
        "sync(core_0, …, core_7)",
        formula_size=9.7,
    )
    add_flow_box(
        slide,
        0.96,
        6.13,
        6.86,
        0.53,
        "目标核本地逐元素归约",
        "Y_i  =  Σ_{k=0}^{7} A_i[k]   [S]",
        title_size=10.9,
        formula_size=8.8,
    )
    add_endpoint(
        slide,
        1.84,
        6.86,
        5.22,
        0.40,
        "结束",
        "每个 Core 得到相同的完整 Y",
        10.8,
    )
    add_legend_panel(
        slide,
        [
            ("i", "当前 / 源 Core 编号"),
            ("j", "目标 Core 编号，j ≠ i"),
            ("S", "完整 tensor 的字节数"),
            ("X_i", "Core_i 持有的完整输入"),
            ("A_j[i]", "Core_j 为源 i 预留的完整 tensor 槽位",),
            ("Y_i", "Core_i 的完整 AllReduce 输出"),
        ],
        [
            "不分片：一次 remote_write 搬运完整 S bytes。",
            "远程写只写独立槽位；sync 后由目标 core 本地求和。",
        ],
    )
    add_common_footer(slide, "AllReduce = 收齐完整输入副本 + 本地逐元素求和；本页未使用 Reduce-Scatter 分片。")


def build_reduce_scatter_slide(prs):
    slide = new_slide(prs)
    add_header(
        slide,
        "8 核 Reduce-Scatter 流程图",
        "每个 core 持有完整 tensor；按目标 core 切成 8 片，归约后每个 core 只保留自己的结果片",
        "2 / 2",
    )
    add_endpoint(
        slide,
        1.56,
        0.95,
        5.72,
        0.53,
        "开始",
        "Core_0 … Core_7 同时启动 Reduce-Scatter",
        11.4,
    )
    for y in (1.55, 2.28, 3.01, 3.81, 4.52, 5.23, 5.94, 6.66):
        add_down_arrow(slide, y)
    add_flow_box(
        slide,
        1.28,
        1.74,
        6.22,
        0.46,
        "每个 Core 准备完整输入",
        "X_i  [S]",
        formula_size=9.8,
    )
    add_flow_box(
        slide,
        0.98,
        2.47,
        6.82,
        0.49,
        "按目标 Core 切成 8 个分片",
        "X_i[0:8]  ·  C = S / 8",
        formula_size=9.0,
    )
    add_flow_box(
        slide,
        1.44,
        3.20,
        5.90,
        0.49,
        "准备本地归约槽位",
        "A_i[0:8]  ·  每槽 C bytes",
        formula_size=8.9,
    )
    add_flow_box(
        slide,
        0.65,
        4.00,
        7.48,
        0.61,
        "把第 j 片写到 Core_j 的对应槽位",
        "A_j[i]  ←  X_i[j]   ·   j=i 本地放置，j≠i remote_write",
        title_size=10.7,
        formula_size=7.9,
    )
    add_flow_box(
        slide,
        1.68,
        4.71,
        5.42,
        0.47,
        "本核写入顺序栅栏",
        "fence core_i",
        formula_size=10.0,
    )
    add_flow_box(
        slide,
        1.42,
        5.42,
        5.94,
        0.51,
        "8 核集合同步",
        "sync(core_0, …, core_7)",
        formula_size=9.7,
    )
    add_flow_box(
        slide,
        0.90,
        6.13,
        6.98,
        0.53,
        "Core_i 本地归约并保留第 i 片",
        "Y_i  =  Σ_{k=0}^{7} A_i[k]   [C]",
        title_size=10.8,
        formula_size=8.8,
    )
    add_endpoint(
        slide,
        1.62,
        6.86,
        5.68,
        0.40,
        "结束",
        "Core_i 持有第 i 个 reduced chunk",
        10.5,
    )
    add_legend_panel(
        slide,
        [
            ("i", "源 / 当前 Core 编号"),
            ("j", "结果 chunk 的目标 Core 编号"),
            ("S", "完整 tensor 的字节数"),
            ("C", "一个 chunk 的字节数，C=S/8"),
            ("X_i[j]", "Core_i 的第 j 个输入分片", 9.0),
            ("A_j[i]", "Core_j 中源 i 的接收槽位", 9.0),
            ("Y_j", "Core_j 最终拥有的第 j 个结果片", 8.9),
        ],
        [
            "Reduce-Scatter：先对同一 chunk 求和，再把结果分散给唯一 owner。",
            "结束后每个 core 只有 1 片；需要 AllGather 才能恢复完整 tensor。",
        ],
    )
    add_common_footer(slide, "Reduce-Scatter = 分片 + 目标核本地归约；输出是分布式的 8 个 reduced chunks。")


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    build_allreduce_slide(prs)
    build_reduce_scatter_slide(prs)
    prs.core_properties.title = "8 核 AllReduce 与 Reduce-Scatter 流程图"
    prs.core_properties.subject = "Remote-write implementation flowcharts for AllReduce and Reduce-Scatter"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Two editable black-and-white flowchart slides with start/end ellipses and symbol definitions."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("allreduce_reducescatter_flowcharts.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()