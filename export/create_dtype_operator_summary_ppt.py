#!/usr/bin/env python3
"""Create a one-slide operator data-type summary for DeepSeek V4 Flash."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from create_moe_route_ppt import (  # noqa: E402
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


ROOT = Path(__file__).resolve().parent.parent
W, H = 13.333, 7.5

PALETTE = {
    "background": "F5F9FA",
    "panel": "FFFFFF",
    "row_alt": "EDF5F6",
    "header": "365E69",
    "ink": "18323A",
    "muted": "55717A",
    "line": "A9C0C7",
    "white": "FFFFFF",
    "cyan": "008C99",
    "blue": "2E72B6",
    "green": "23855F",
    "amber": "B76700",
    "coral": "C6544C",
    "purple": "7050A2",
}


def color(name: str) -> str:
    return PALETTE.get(name, name)


def set_cell_border(cell, line_color: str = "line", line_width: float = 0.55) -> None:
    cell_properties = cell._tc.get_or_add_tcPr()
    width = str(int(line_width * 12700))
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        edge_tag = qn(edge)
        line = cell_properties.find(edge_tag)
        if line is None:
            line = OxmlElement(edge)
            cell_properties.append(line)
        line.set("w", width)
        line.set("cap", "flat")
        line.set("cmpd", "sng")
        line.set("algn", "ctr")
        solid_fill = line.find(qn("a:solidFill"))
        if solid_fill is None:
            solid_fill = OxmlElement("a:solidFill")
            line.append(solid_fill)
        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = OxmlElement("a:srgbClr")
            solid_fill.append(srgb)
        srgb.set("val", color(line_color))
        preset_dash = line.find(qn("a:prstDash"))
        if preset_dash is None:
            preset_dash = OxmlElement("a:prstDash")
            line.append(preset_dash)
        preset_dash.set("val", "solid")


def set_cell_text(
    cell,
    text: str,
    fill: str = "panel",
    text_color: str = "ink",
    size: float = 6.2,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    font: str = FONT_CJK,
    line_color: str = "line",
    line_width: float = 0.55,
) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(color(fill))
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.025)
    cell.margin_bottom = Inches(0.025)
    cell.vertical_anchor = valign
    cell.text = text
    text_frame = cell.text_frame
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.space_before = 0
        paragraph.space_after = 0
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color(text_color))
    set_cell_border(cell, line_color, line_width)


def add_native_table(slide, x: float, y: float, widths: list[float], headers: list[str], groups) -> None:
    header_h, group_h, row_h = 0.36, 0.20, 0.36
    row_count = 1 + sum(1 + len(rows) for _, _, _, rows in groups)
    table_h = header_h + sum(
        group_h + len(rows) * row_h for _, _, _, rows in groups
    )
    graphic_frame = slide.shapes.add_table(
        row_count,
        len(widths),
        Inches(x),
        Inches(y),
        Inches(sum(widths)),
        Inches(table_h),
    )
    table = graphic_frame.table
    for column, width in zip(table.columns, widths):
        column.width = Inches(width)

    table.rows[0].height = Inches(header_h)
    for column, header in enumerate(headers):
        set_cell_text(
            table.cell(0, column),
            header,
            fill="header",
            text_color="white",
            size=6.9,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            font=FONT,
            line_width=0.65,
        )

    row_index = 1
    alternate = False
    for title, detail, accent, rows in groups:
        table.rows[row_index].height = Inches(group_h)
        set_cell_text(
            table.cell(row_index, 0),
            title,
            fill="row_alt",
            text_color=accent,
            size=6.8,
            bold=True,
            font=FONT,
            line_width=0.55,
        )
        detail_cell = table.cell(row_index, 1)
        detail_cell.merge(table.cell(row_index, len(widths) - 1))
        set_cell_text(
            detail_cell,
            detail,
            fill="row_alt",
            text_color="muted",
            size=6.1,
            align=PP_ALIGN.RIGHT,
            font=FONT_CJK,
            line_width=0.55,
        )
        row_index += 1
        for values in rows:
            table.rows[row_index].height = Inches(row_h)
            fill = "row_alt" if alternate else "panel"
            for column, value in enumerate(values):
                set_cell_text(
                    table.cell(row_index, column),
                    value,
                    fill=fill,
                    text_color="ink",
                    size=6.2 if column else 6.55,
                    bold=column == 0,
                    valign=MSO_ANCHOR.MIDDLE,
                    font=FONT_CJK,
                )
            alternate = not alternate
            row_index += 1


def build_deck(output: Path) -> None:
    config = json.loads((ROOT / "config.json").read_text(encoding="utf-8"))
    quantization = config["quantization_config"]
    config_line = (
        f"linear/shared: {quantization['quant_method'].upper()} {quantization['fmt'].upper()} + {quantization['scale_fmt'].upper()}\n"
        f"routed: {config['expert_dtype'].upper()} E2M1 + {quantization['scale_fmt'].upper()} · state: BF16"
    )

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PALETTE["background"])

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.045, PALETTE["cyan"])
    add_text(
        slide,
        "DeepSeek V4 Flash · Operator Data Types",
        0.38,
        0.15,
        8.15,
        0.38,
        21.0,
        color("ink"),
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "算子设计总览：权重存储、量化 Scale、激活 / hidden state 与 TP 组织方式",
        0.41,
        0.59,
        8.65,
        0.22,
        8.9,
        color("muted"),
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.38, 0.16, 3.56, 0.36, color("row_alt"), color("cyan"), 0.75)
    add_text(
        slide,
        config_line,
        9.49,
        0.19,
        3.34,
        0.30,
        6.2,
        color("cyan"),
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.38, 0.91, 12.56, 0.36, color("panel"), color("line"), 0.75)
    legend = [
        ("FP8 E4M3  weight", "blue", 2.05),
        ("FP4 E2M1  routed expert", "green", 2.35),
        ("UE8M0 / E8M0  Scale", "amber", 2.65),
        ("BF16  state / cache", "cyan", 2.25),
        ("FP32  compute", "purple", 1.80),
    ]
    legend_x = 0.58
    for label, accent, width in legend:
        add_text(
            slide,
            label,
            legend_x,
            1.00,
            width,
            0.17,
            6.8,
            color(accent),
            True,
            CODE_FONT,
            PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )
        legend_x += width + 0.10

    table_x, table_y = 0.38, 1.45
    widths = [2.35, 2.62, 2.05, 2.72, 2.82]
    table_w = sum(widths)
    headers = [
        "OPERATOR / TENSOR",
        "WEIGHT / STORAGE",
        "SCALE TYPE",
        "ACTIVATION / HIDDEN STATE",
        "DESIGN / PARALLELISM",
    ]
    groups = [
        (
            "CORE STATE",
            "hidden state 是 BF16 主路径",
            "cyan",
            [
                (
                    "Hidden state / residual",
                    "BF16 hidden\nBF16 residual",
                    "N/A",
                    "BF16 state tensors;\nFP32 HC mixing when enabled",
                    "完整 [B,S,D] 视图；mHC streams\n在 rank 内复制，不切 TP",
                )
            ],
        ),
        (
            "ATTENTION",
            "主投影、KV、Indexer 与 Compressor",
            "blue",
            [
                (
                    "Main projections\nwq_a / wq_b / wkv / wo_b",
                    "FP8 E4M3\nweights",
                    "UE8M0 (E8M0)\nunsigned exponent-only",
                    "FP8 GEMM; BF16 input\ndynamically quantized",
                    "wq_a / wkv replicate;\nQ heads 与 O groups 按 TP 切分",
                ),
                (
                    "Output projection\nwo_a",
                    "BF16 inference\nweights",
                    "N/A",
                    "BF16 tensor product\n(einsum)",
                    "checkpoint FP8 tensor 在推理实现中\n转换并按 BF16 使用",
                ),
                (
                    "Main KV cache",
                    "BF16 unquantized\ncache entries",
                    "N/A\n(unquantized cache)",
                    "BF16 cache\nraw + compressed entries",
                    "window / ratio-4 / ratio-128\nKV 在每个 TP Rank 复制",
                ),
                (
                    "Indexer KV + projections",
                    "KV: BF16 unquantized\nwq_b: FP8; weights_proj: BF16",
                    "wq_b: UE8M0\nKV / weights_proj: N/A",
                    "QAT/FP4 simulated QKV;\nFP32 scoring",
                    "仅 ratio-4 层；Indexer KV\n与权重按当前实现组织",
                ),
                (
                    "Main Compressor\nwkv / wgate / ape",
                    "checkpoint mainly BF16;\nFP32 inference parameters",
                    "N/A",
                    "FP32 compression /\nsoftmax",
                    "生成压缩 KV；ratio-4 /\nratio-128 层使用",
                ),
            ],
        ),
        (
            "MoE",
            "专家权重与路由数据类型",
            "green",
            [
                (
                    "Routed expert weights\nw1 / w2 / w3",
                    "FP4 E2M1\npacked weights",
                    "UE8M0 (E8M0)\nunsigned scale",
                    "FP4 GEMM; SwiGLU FP32;\noutput cast to BF16",
                    "expert_id 分片；每 token\nTop-K=6 个路由专家",
                ),
                (
                    "Shared expert weights\nw1 / w2 / w3",
                    "FP8 E4M3\nweights",
                    "UE8M0 (E8M0)\nunsigned scale",
                    "FP8 GEMM;\nSwiGLU FP32",
                    "每层 1 个共享专家；\n跨 TP Rank 复制",
                ),
                (
                    "Router / Token-ID table",
                    "FP32 bias;\nINT32 lookup table",
                    "N/A",
                    "FP32 score, Top-K\nand normalization",
                    "前 3 层 Token-ID 查表；\n其余层 FP32 score routing",
                ),
            ],
        ),
        (
            "NORMALIZATION / TAIL",
            "数值稳定路径与词表输出",
            "amber",
            [
                (
                    "RMSNorm / mHC / sinks",
                    "BF16 weights;\nFP32 implementation params",
                    "N/A",
                    "FP32 compute →\nBF16 output / residual",
                    "Norm、Sinkhorn、HC mixing\n保持 rank 内完整视图",
                ),
                (
                    "LM head / logits",
                    "checkpoint BF16;\nFP32 inference parameters",
                    "N/A",
                    "FP32 linear + logits;\nBF16 hidden input",
                    "vocab-sharded；local logits\n最后 AllGather 拼接",
                ),
            ],
        ),
    ]

    add_native_table(slide, table_x, table_y, widths, headers, groups)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.38, 7.10, 12.56, 0.23, color("row_alt"), color("line"), 0.65)
    add_text(
        slide,
        "口径：Scale 只对应量化权重；hidden state / residual / KV cache 均按 BF16 存储。wo_a 是推理 BF16，不是“FP8 + Scale”。",
        0.55,
        7.145,
        11.95,
        0.15,
        6.7,
        color("ink"),
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "source: config.json · dtype sheet · current inference design",
        8.58,
        7.36,
        4.20,
        0.12,
        5.8,
        color("muted"),
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "DeepSeek V4 Flash Operator Data Types"
    prs.core_properties.subject = "Weights, scales, hidden states, KV cache and operator design"
    prs.core_properties.author = "DeepSeek V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide operator data-type summary generated from the reference configuration."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("dtype_operator_summary.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()