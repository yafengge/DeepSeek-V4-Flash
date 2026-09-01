#!/usr/bin/env python3
"""Create a two-slide editable Attention parameter and compute summary."""

from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "calculate"))

from generate_calculator import (  # noqa: E402
    LayerCounts,
    cache_values,
    count_layers,
    load_inputs,
    parameter_components,
    replace,
    scenario_items,
    summarize_items,
)
from create_moe_route_ppt import (  # noqa: E402
    CODE_FONT,
    COLORS,
    FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5
TP_VALUES = (1, 8)


@dataclass(frozen=True)
class SectionSpec:
    title: str
    subtitle: str
    config: object
    layers: LayerCounts
    accent: str


@dataclass(frozen=True)
class RowData:
    tp: int
    params: str
    cache: str
    dtype: str
    compute: str
    note: str


def format_count(value: float) -> str:
    if abs(value) >= 1e9:
        return f"{value / 1e9:.3f}B"
    if abs(value) >= 1e6:
        return f"{value / 1e6:.3f}M"
    if abs(value) >= 1e3:
        return f"{value / 1e3:.3f}K"
    return f"{value:.0f}"


def format_gb(value: float) -> str:
    return f"{value / 1e9:.3f} GB"


def format_flops(value: float, mode: str) -> str:
    if mode == "prefill":
        return f"{value / 1e12:.3f} TFLOPs"
    return f"{value / 1e9:.3f} GFLOPs"


def component_map(config, layers: LayerCounts) -> dict[str, object]:
    return {
        component.name: component
        for component in parameter_components(config, layers)
        if component.category == "Attention"
    }


def parameter_text(config, layers: LayerCounts) -> str:
    components = component_map(config, layers)
    core = components["Core Q/K/O projections"]
    compressors = components["KV compressors"]
    indexers = components["Ratio-4 Indexers"]
    total_count = sum(component.rank_count for component in components.values())
    total_bytes = sum(component.rank_bytes for component in components.values())
    return (
        f"合计 {format_count(total_count)} logical\n"
        f"Core {format_count(core.rank_count)} · Comp {format_count(compressors.rank_count)}\n"
        f"Indexer {format_count(indexers.rank_count)}\n"
        f"存储 {format_gb(total_bytes)}"
    )


def cache_text(config, layers: LayerCounts, context: int) -> str:
    cache = cache_values(config, layers, context, 1)
    auxiliary = cache["indexer"] + cache["states"]
    return (
        f"主 KV {format_gb(cache['main'])}\n"
        f"Indexer {format_gb(cache['indexer'])}\n"
        f"State {format_gb(cache['states'])}\n"
        f"合计 {format_gb(cache['total'])}"
    )


def dtype_text(section_key: str) -> str:
    if section_key == "ratio4":
        return (
            "主投影：FP8 E4M3 + E8M0\n"
            "wo_a / Q / KV / attn：BF16\n"
            "Norm / Sink / Compressor：FP32\n"
            "Indexer：FP8 / BF16 / FP32"
        )
    if section_key == "ratio8":
        return (
            "主投影：FP8 E4M3 + E8M0\n"
            "wo_a / Q / KV / attn：BF16\n"
            "Norm / Compressor：FP32\n"
            "Indexer：无"
        )
    return (
        "主投影：FP8 E4M3 + E8M0\n"
        "wo_a / Q / KV / attn：BF16\n"
        "Norm / Sink：FP32\n"
        "含 ratio-4 Indexer：FP8 / BF16 / FP32"
    )


def compute_text(config, layers: LayerCounts, mode: str) -> str:
    items, _ = scenario_items(config, layers, mode)
    attention_items = [item for item in items if item.category == "Attention"]
    by_name = {item.name: item.rank_flops for item in attention_items}
    total = summarize_items(items)["attention_major_flops_per_rank"]
    projection = by_name["Q/K/O projections"]
    sparse = by_name["Sparse attention QK + AV"]
    indexer = (
        by_name["Ratio-4 Indexer projections"]
        + by_name["Ratio-4 Indexer score scan"]
        + by_name["Indexer compressor projections"]
    )
    compressor = by_name["Main KV compressor projections"]
    lines = [
        f"总计 {format_flops(total, mode)}",
        f"Proj {format_flops(projection, mode)} · QK+AV {format_flops(sparse, mode)}",
    ]
    if indexer:
        lines.append(f"Indexer {format_flops(indexer, mode)}")
    if compressor:
        lines.append(f"Compressor {format_flops(compressor, mode)}")
    return "\n".join(lines)


def make_sections(config, mode: str) -> list[tuple[str, SectionSpec]]:
    native_layers = count_layers(config[1], config[0])
    base = config[0]
    return [
        (
            "full",
            SectionSpec(
                "整网 Attention",
                "原生 43 层：window×2 + ratio=4×21 + ratio=128×20",
                base,
                native_layers,
                "cyan",
            ),
        ),
        (
            "ratio4",
            SectionSpec(
                "ratio=4 Attention",
                "单个短压缩层：含 ratio-4 Indexer 与 Top-K=512",
                base,
                LayerCounts(total=1, window=0, short=1, long=0, hash=0),
                "blue",
            ),
        ),
        (
            "ratio8",
            SectionSpec(
                "ratio=8 Attention",
                "单个长压缩层：公式对比场景，原生模型未启用 ratio=8",
                replace(base, long_ratio=8),
                LayerCounts(total=1, window=0, short=0, long=1, hash=0),
                "amber",
            ),
        ),
    ]


def add_table_cell(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    fill: str = "panel",
    color: str = "black",
    size: float = 7.0,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font=FONT_CJK,
):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, "line", 0.65)
    return add_text(
        slide,
        text,
        x + 0.08,
        y + 0.06,
        w - 0.16,
        h - 0.12,
        size,
        color,
        bold,
        font,
        align,
        valign,
    )


def add_section_band(slide, x: float, y: float, w: float, h: float, key: str, section: SectionSpec):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel_2", "line", 0.7)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.075, h, section.accent)
    add_text(
        slide,
        section.title,
        x + 0.16,
        y + 0.035,
        2.15,
        h - 0.07,
        9.4,
        section.accent,
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        section.subtitle,
        x + 2.35,
        y + 0.045,
        w - 2.52,
        h - 0.09,
        7.0,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )


def build_rows(section_key: str, section: SectionSpec, mode: str, context: int) -> list[RowData]:
    rows = []
    for tp in TP_VALUES:
        config = replace(section.config, tp=tp)
        local_heads = int(config.heads / tp)
        local_groups = int(config.o_groups / tp)
        rows.append(
            RowData(
                tp=tp,
                params=parameter_text(config, section.layers),
                cache=cache_text(config, section.layers, context),
                dtype=dtype_text(section_key),
                compute=compute_text(config, section.layers, mode),
                note=f"Q {local_heads} heads\nO {local_groups} groups\n每卡 / rank",
            )
        )
    return rows


def add_summary_table(slide, mode: str, config) -> None:
    table_x, table_y, table_w = 0.28, 0.96, 12.78
    header_h, section_h, row_h = 0.40, 0.30, 0.72
    column_widths = [0.82, 3.60, 1.62, 2.54, 2.50, 1.70]
    context = config[0].prefill_sequence if mode == "prefill" else config[0].decode_context
    sections = make_sections(config, mode)
    table_h = header_h + len(sections) * (section_h + 2 * row_h)

    add_shape(slide, MSO_SHAPE.RECTANGLE, table_x, table_y, table_w, table_h, "panel", "line", 0.9)
    cursor_x = table_x
    headers = ["TP / 每卡", "Attention 算子参数", "KV Cache 容量", "数据类型", "计算量 / 卡", "分片状态"]
    for width, header in zip(column_widths, headers):
        add_table_cell(
            slide,
            cursor_x,
            table_y,
            width,
            header_h,
            header,
            fill="white",
            color="muted",
            size=7.4,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        cursor_x += width

    y = table_y + header_h
    for section_key, section in sections:
        add_section_band(slide, table_x, y, table_w, section_h, section_key, section)
        y += section_h
        for row in build_rows(section_key, section, mode, context):
            cursor_x = table_x
            tp_color = "cyan" if row.tp == 1 else "blue"
            add_table_cell(
                slide,
                cursor_x,
                y,
                column_widths[0],
                row_h,
                f"TP{row.tp}\n每卡",
                fill="panel_2",
                color=tp_color,
                size=9.0,
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                font=FONT,
            )
            cursor_x += column_widths[0]
            add_table_cell(slide, cursor_x, y, column_widths[1], row_h, row.params, size=6.7)
            cursor_x += column_widths[1]
            add_table_cell(slide, cursor_x, y, column_widths[2], row_h, row.cache, size=6.7)
            cursor_x += column_widths[2]
            add_table_cell(slide, cursor_x, y, column_widths[3], row_h, row.dtype, size=6.55)
            cursor_x += column_widths[3]
            add_table_cell(slide, cursor_x, y, column_widths[4], row_h, row.compute, size=6.55)
            cursor_x += column_widths[4]
            add_table_cell(
                slide,
                cursor_x,
                y,
                column_widths[5],
                row_h,
                row.note,
                size=7.0,
                color="muted",
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
            )
            y += row_h


def build_slide(prs: Presentation, mode: str, config) -> None:
    is_prefill = mode == "prefill"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.045, "cyan")
    add_text(
        slide,
        f"Attention 参数、KV Cache 与算力汇总 · {'Prefill' if is_prefill else 'Decode'}",
        0.34,
        0.12,
        9.75,
        0.34,
        19.0,
        "red",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    scenario_label = "B=1 · S=8,192" if is_prefill else "B=1 · 1 token · C=1,048,576"
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.38, 0.13, 2.57, 0.32, "panel_2", "blue", 0.8)
    add_text(
        slide,
        scenario_label,
        10.48,
        0.20,
        2.37,
        0.17,
        7.3,
        "blue",
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "三段同一张表：整网原生层级 / ratio=4 单层 / ratio=8 对比单层；每段垂直拆分 TP1、TP8 每卡状态",
        0.37,
        0.56,
        9.65,
        0.19,
        8.0,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "D=4096 · H=64 · d=512 · q=1024 · G=8",
        10.03,
        0.56,
        2.92,
        0.19,
        7.2,
        "muted",
        True,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    add_summary_table(slide, mode, config)
    add_text(
        slide,
        "KV Cache 为有效驻留容量：主 KV + Indexer cache + Compressor state；GB=10^9 bytes。ratio=8 是公式对比假设，原生长压缩配置为 ratio=128。",
        0.36,
        6.83,
        12.45,
        0.23,
        6.7,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "MAC=2 FLOPs · Attention major 口径 · source: inference/model.py / calculate/generate_calculator.py",
        0.36,
        7.16,
        12.45,
        0.17,
        6.6,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )


def build_deck(output: Path) -> None:
    config = load_inputs(ROOT / "inference" / "config.json")
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    build_slide(prs, "prefill", config)
    build_slide(prs, "decode", config)
    prs.core_properties.title = "DeepSeek V4 Flash Attention Prefill Decode Summary"
    prs.core_properties.subject = "TP1/TP8 per-card Attention parameters, KV cache, dtype, and FLOPs"
    prs.core_properties.author = "DeepSeek V4 Flash repository"
    prs.core_properties.comments = "Editable two-slide summary generated from the repository calculator."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("attention_prefill_decode_summary.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()