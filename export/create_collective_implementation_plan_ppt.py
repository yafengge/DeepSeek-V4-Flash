#!/usr/bin/env python3
"""Create a one-slide per-collective implementation plan."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from create_moe_route_ppt import (
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5
TABLE_X, TABLE_Y, TABLE_W = 0.42, 1.18, 12.49
COLUMNS = [1.32, 2.44, 2.38, 2.57, 1.82, 1.96]


def add_cell(slide, x, y, w, h, title, detail, accent=None, align=PP_ALIGN.LEFT):
    fill = "panel" if int((y - TABLE_Y) * 100) % 2 else "light"
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, "line", 0.75)
    if accent:
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.055, h, accent)
    text_x = x + 0.12
    text_w = w - 0.20
    add_text(
        slide,
        title,
        text_x,
        y + 0.08,
        text_w,
        0.25,
        8.0 if w > 1.7 else 7.5,
        "black",
        True,
        FONT_CJK,
        align,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        text_x,
        y + 0.35,
        text_w,
        h - 0.42,
        6.8 if w > 1.7 else 6.5,
        "black",
        False,
        CODE_FONT if "remote" in detail or "NoC" in detail else FONT_CJK,
        align,
        MSO_ANCHOR.TOP,
    )


def add_operation_cell(slide, x, y, w, h, title, detail, accent):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel", accent, 1.4)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, accent)
    add_text(
        slide,
        title,
        x + 0.16,
        y + 0.16,
        w - 0.26,
        0.28,
        9.2,
        "black",
        True,
        FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        x + 0.14,
        y + 0.49,
        w - 0.24,
        h - 0.57,
        7.2,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.TOP,
    )


def add_header_cell(slide, x, y, w, h, text):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel_2", "black", 0.9)
    add_text(
        slide,
        text,
        x + 0.06,
        y + 0.04,
        w - 0.12,
        h - 0.08,
        7.4,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_info_box(slide, x, y, w, h, title, detail, accent):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "light", accent, 1.0)
    add_text(
        slide,
        title,
        x + 0.14,
        y + 0.08,
        w - 0.28,
        0.22,
        8.6,
        accent,
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        x + 0.14,
        y + 0.34,
        w - 0.28,
        h - 0.41,
        7.3,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_text(
        slide,
        "集合通信实现方案（逐算子）",
        0.42,
        0.16,
        7.25,
        0.39,
        21.0,
        "red",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "设计目标：Broadcast 由 NoC multicast descriptor 支持；其余通信由 Core kernel 完成数据搬运后的拼接或求和",
        0.45,
        0.61,
        9.45,
        0.22,
        8.8,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "NoC 传输 / Core 计算",
        10.15,
        0.25,
        2.72,
        0.22,
        9.2,
        "black",
        True,
        FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "主路径：P2P / multicast + local kernel + event/fence",
        8.15,
        0.66,
        4.72,
        0.19,
        7.4,
        "black",
        True,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    header_y, header_h = TABLE_Y, 0.46
    cursor = TABLE_X
    headers = [
        "集合通信 / OP",
        "推荐路径",
        "Core kernel / 计算",
        "搬运 / NoC 路径",
        "缓存与同步",
        "代价 / 备注",
    ]
    for width, header in zip(COLUMNS, headers):
        add_header_cell(slide, cursor, header_y, width, header_h, header)
        cursor += width

    rows = [
        (
            0.84,
            "Broadcast",
            "广播",
            "cyan",
            [
                ("NoC multicast descriptor", "1 个 BCAST；源端只读 1 次"),
                ("不做数值归约", "各 Core 直接消费本地 B_i"),
                ("multicast tree", "NoC 在分支复制 packet/flit；目标写本地 SRAM/DRAM"),
                ("completion bitmap", "fence + sync"),
                ("面积低", "需组播；无需 Reduce Engine"),
            ],
        ),
        (
            0.91,
            "AllGather",
            "全收集",
            "blue",
            [
                ("Core kernel + remote SRAM", "remote DRAM copy = 容量兜底"),
                ("放置本地分片", "Y_i[i] ← X_i；只拼接，不相加"),
                ("P2P remote copy", "每个 X_i 写入其他 Core 的 slot；优先 remote SRAM"),
                ("slot ready", "fence + sync"),
                ("实现简单", "SRAM 不足时落 DRAM；通信约为 (P−1)S/P"),
            ],
        ),
        (
            0.98,
            "Reduce-Scatter",
            "规约分发",
            "green",
            [
                ("Core kernel + vector stream-in", "P2P / Ring；按 chunk 流水"),
                ("本地向量加法", "A_i[c] += RX[c]；Core vector 逐 chunk 累加"),
                ("NoC → Rx FIFO", "不落远程 DRAM；双缓冲 + credit/backpressure"),
                ("ready / consumed", "chunk_id + step"),
                ("无 Reduce Engine", "占用 Core vector 与本地 SRAM 带宽"),
            ],
        ),
        (
            1.05,
            "AllReduce",
            "全规约",
            "purple",
            [
                ("RS + AG 组合", "RS: vector stream-in；AG: remote SRAM copy"),
                ("两阶段 kernel", "先 chunk add，再交换归约后的 chunk"),
                ("P2P + remote copy", "先规约分片，再全收集；不在 Router 求和"),
                ("phase fence", "completion counter + final sync"),
                ("硬件省", "Core/NoC 流量与总延迟高于原生归约"),
            ],
        ),
    ]

    row_y = TABLE_Y + header_h
    for row_h, operation, operation_cn, accent, cells in rows:
        add_operation_cell(slide, TABLE_X, row_y, COLUMNS[0], row_h, operation, operation_cn, accent)
        cursor = TABLE_X + COLUMNS[0]
        for column_index, (title, detail) in enumerate(cells, start=1):
            add_cell(
                slide,
                cursor,
                row_y,
                COLUMNS[column_index],
                row_h,
                title,
                detail,
                align=PP_ALIGN.LEFT,
            )
            cursor += COLUMNS[column_index]
        row_y += row_h

    add_shape(slide, MSO_SHAPE.RECTANGLE, TABLE_X, row_y, TABLE_W, 0.018, "black")

    add_text(
        slide,
        "实现边界",
        0.45,
        5.58,
        1.05,
        0.20,
        8.2,
        "gray",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "NoC 原生 collective：仅保留 BCAST；Core-assisted collective：NoC 负责 P2P/组播传输，Core 负责 kernel、vector add 和 accumulator。",
        1.52,
        5.58,
        11.28,
        0.20,
        7.9,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_info_box(
        slide,
        0.45,
        6.02,
        6.08,
        0.91,
        "路径选择",
        "性能倾向：remote vector stream-in > remote SRAM copy > remote DRAM copy\n前提：Core vector 空闲，且本地 SRAM 带宽能够承接流入数据。",
        "blue",
    )
    add_info_box(
        slide,
        6.78,
        6.02,
        6.10,
        0.91,
        "最小硬件接口",
        "NoC：P2P + multicast + FIFO/credit + event；Core：kernel + vector add + accumulator\n同步：ready/consumed + completion + fence；Broadcast 不需要 Reduce Engine。",
        "green",
    )

    prs.core_properties.title = "集合通信逐算子实现方案"
    prs.core_properties.subject = "Broadcast via NoC multicast and Core-assisted collective communication paths"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide implementation plan for Broadcast, AllGather, Reduce-Scatter and AllReduce."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("collective_implementation_plan.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()