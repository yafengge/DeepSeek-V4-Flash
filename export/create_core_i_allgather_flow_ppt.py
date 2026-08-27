#!/usr/bin/env python3
"""Create a one-slide top-down Core_i remote-write AllGather flowchart."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from create_moe_route_ppt import (
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_connector,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5


def add_flow_node(
    slide,
    number,
    y,
    h,
    title,
    formula,
    detail,
    accent,
    formula_size=9.2,
):
    x, w = 0.45, 7.05
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "panel", "line", 0.9)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.075, h, accent)
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.18, y + 0.07, 0.30, 0.30, accent)
    add_text(
        slide,
        str(number),
        x + 0.18,
        y + 0.07,
        0.30,
        0.30,
        8.5,
        "bg",
        True,
        FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        title,
        x + 0.60,
        y + 0.06,
        2.54,
        0.29,
        10.8,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        formula,
        x + 3.15,
        y + 0.06,
        3.67,
        0.29,
        formula_size,
        accent,
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        x + 0.60,
        y + 0.36,
        6.22,
        h - 0.39,
        7.3,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def add_down_arrow(slide, y, color="line"):
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, 3.83, y, 0.22, 0.13, color)


def add_core_chip(slide, label, x, y, accent):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.82, 0.27, "panel", accent, 0.75)
    add_text(
        slide,
        label,
        x + 0.03,
        y + 0.03,
        0.76,
        0.21,
        7.0,
        accent,
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_target_row(slide, label, y, accent="cyan", muted=False):
    fill = "light" if muted else "panel_2"
    line = "line" if muted else accent
    color = "muted" if muted else "black"
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.84, y, 2.73, 0.36, fill, line, 0.75)
    add_text(
        slide,
        label,
        9.98,
        y + 0.05,
        2.45,
        0.25,
        8.0,
        color,
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.045, "cyan")
    add_text(
        slide,
        "8 核远程写入 AllGather：Core_i 单核执行流程",
        0.42,
        0.15,
        10.30,
        0.43,
        21.0,
        "red",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "八个核心同时启动同一个算子；本图只展开任意一个 Core_i 的控制流与远程内存动作",
        0.45,
        0.62,
        12.20,
        0.20,
        9.1,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.42, 0.94, 12.46, 0.42, "panel_2", "line", 0.8)
    add_text(
        slide,
        "SIMULTANEOUS LAUNCH",
        0.62,
        1.04,
        1.75,
        0.20,
        7.6,
        "black",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    chip_x = 2.55
    chip_colors = ["cyan", "blue", "green", "amber", "coral", "purple", "cyan", "blue"]
    for core, accent in enumerate(chip_colors):
        add_core_chip(slide, f"core_{core}", chip_x + core * 0.91, 1.02, accent)
    add_text(
        slide,
        "launch AG_i",
        10.03,
        1.04,
        1.08,
        0.20,
        8.0,
        "blue",
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "i ∈ {0…7}",
        11.36,
        1.04,
        1.17,
        0.20,
        8.0,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    add_text(
        slide,
        "CORE_i / LOCAL CONTROL FLOW",
        0.45,
        1.51,
        7.05,
        0.16,
        7.5,
        "muted",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "REMOTE MEMORY VIEW",
        7.75,
        1.51,
        5.13,
        0.16,
        7.5,
        "muted",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    nodes = [
        (1.70, 0.44, "启动 AllGather", "Core_i  ·  X_i [C]", "八个 core 同时进入同一算子", "cyan", 9.0),
        (2.29, 0.48, "准备接收缓冲区", "Y_i[0:8]  ·  ready_i=0", "每个 slot 预留 C bytes", "blue", 8.6),
        (2.92, 0.44, "放置本地分片", "Y_i[i] ← X_i", "local copy / alias；不走 NoC", "green", 9.2),
        (3.51, 0.50, "生成远程写描述符", "dst = Y_j + i·C  (j ≠ i)", "目标 core j 的第 i 个接收槽位", "amber", 8.6),
        (4.16, 0.58, "提交 7 个远程写", "remote_write(X_i → Y_j[i], C) × 7", "DMA / NoC；可并发 outstanding；不做加法", "cyan", 8.5),
        (4.90, 0.48, "本核写入顺序栅栏", "fence core_i", "payload writes 先完成并对后续同步可见", "coral", 10.2),
        (5.53, 0.52, "8 核集合同步", "sync core_0~8", "所有参与者到达 fence 后再继续", "purple", 10.0),
        (6.17, 0.48, "读取完整结果", "Y_i = [X_0, X_1, …, X_7]", "全部 slot 有效；交给后续算子", "green", 8.8),
    ]
    for number, node in enumerate(nodes, start=1):
        add_flow_node(slide, number, *node)
    for y in (2.18, 2.82, 3.42, 4.08, 4.78, 5.40, 6.05):
        add_down_arrow(slide, y)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.75, 1.70, 5.13, 5.02, "panel", "line", 0.9)
    add_text(
        slide,
        "Core_i 的一个源分片",
        7.98,
        1.91,
        2.02,
        0.22,
        10.0,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.98, 2.22, 1.48, 0.61, "panel_2", "cyan", 1.0)
    add_text(
        slide,
        "X_i",
        8.12,
        2.30,
        1.20,
        0.22,
        12.0,
        "cyan",
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "[C bytes]",
        8.11,
        2.57,
        1.22,
        0.16,
        7.2,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "写入同一个 offset：i·C",
        9.82,
        2.03,
        2.54,
        0.20,
        8.0,
        "amber",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_connector(slide, 9.48, 2.52, 9.78, 2.52, "cyan", 1.6)
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, 9.61, 2.45, 0.20, 0.14, "cyan")
    add_connector(slide, 9.78, 2.52, 9.78, 6.12, "cyan", 1.1)
    target_rows = [
        ("core_0 : Y_0[i]", 2.43, False),
        ("core_1 : Y_1[i]", 2.87, False),
        ("core_2 : Y_2[i]", 3.31, False),
        ("core_3 : Y_3[i]", 3.75, False),
        ("...     : ...", 4.19, True),
        ("core_5 : Y_5[i]", 4.63, False),
        ("core_6 : Y_6[i]", 5.07, False),
        ("core_7 : Y_7[i]", 5.51, False),
    ]
    for label, y, muted in target_rows:
        add_connector(slide, 9.78, y + 0.18, 9.84, y + 0.18, "line" if muted else "cyan", 0.9)
        add_target_row(slide, label, y, "cyan", muted)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.98, 6.05, 4.58, 0.40, "panel_2", "green", 0.9)
    add_text(
        slide,
        "Y_j[i] ← X_i  ·  j ≠ i：7 次 remote_write",
        8.15,
        6.15,
        4.24,
        0.20,
        8.0,
        "green",
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "j=i 的 Y_i[i] 已在步骤 ③ 本地放置",
        8.00,
        6.51,
        4.55,
        0.17,
        7.4,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.45, 6.86, 12.43, 0.39, "panel_2", "line", 0.8)
    add_text(
        slide,
        "同步分工",
        0.65,
        6.96,
        0.78,
        0.18,
        8.0,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "remote_write = 搬运",
        1.63,
        6.96,
        1.70,
        0.18,
        8.0,
        "cyan",
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "fence core_i = 顺序",
        3.53,
        6.96,
        1.70,
        0.18,
        8.0,
        "coral",
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "sync core_0~8 = 集合相遇",
        5.45,
        6.96,
        2.54,
        0.18,
        8.0,
        "purple",
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "8 核编号通常为 core_0…core_7",
        8.20,
        6.96,
        2.46,
        0.18,
        7.5,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "每核 7C 发 + 7C 收；无 Reduce Engine",
        10.85,
        6.96,
        1.79,
        0.18,
        7.2,
        "amber",
        True,
        FONT_CJK,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Remote-write AllGather  ·  Core_i perspective  ·  editable one-slide diagram",
        0.45,
        7.31,
        7.80,
        0.14,
        6.7,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "X_i → Y_j[i]  ·  j ≠ i",
        10.40,
        7.31,
        2.48,
        0.14,
        6.7,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "8 核远程写入 AllGather：Core_i 单核执行流程"
    prs.core_properties.subject = "Top-down single-core remote-memory-write AllGather flowchart"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide diagram showing the Core_i control flow, remote writes, fence, and eight-core synchronization."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("core_i_remote_write_allgather_flow.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()