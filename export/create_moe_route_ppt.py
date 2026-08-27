#!/usr/bin/env python3
"""Create a one-slide editable flowchart for the DeepSeek-V4 MoE router."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches, Pt


W, H = 13.333, 7.5
FONT = "Arial"
FONT_CJK = "Microsoft YaHei"
CODE_FONT = "Consolas"

COLORS = {
    "bg": "F6F9FB",
    "panel": "FFFFFF",
    "panel_2": "EAF2F5",
    "light": "F3F3F3",
    "red": "C00000",
    "black": "111111",
    "gray": "555555",
    "line": "C2D4DE",
    "muted": "5B6F7D",
    "white": "142C3A",
    "cyan": "008C99",
    "blue": "2E72B6",
    "green": "23855F",
    "amber": "B76700",
    "coral": "C6544C",
    "purple": "7050A2",
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS.get(fill, fill))
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(COLORS.get(line, line))
        shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size,
    color="white",
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.background()
    box.line.fill.background()
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(margin)
    frame.margin_right = Pt(margin)
    frame.margin_top = Pt(margin)
    frame.margin_bottom = Pt(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(COLORS.get(color, color))
    return box


def add_pill(slide, text, x, y, w, h, accent):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "panel_2", accent, 0.8)
    add_text(
        slide,
        text,
        x + 0.08,
        y + 0.02,
        w - 0.16,
        h - 0.04,
        8.3,
        accent,
        True,
        FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_arrow(slide, x, y, w, h, color="line", direction="right"):
    shape_type = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }[direction]
    arrow = add_shape(slide, shape_type, x, y, w, h, color)
    return arrow


def add_connector(slide, x1, y1, x2, y2, color="line", width=1.4, dashed=False):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(COLORS.get(color, color))
    connector.line.width = Pt(width)
    if dashed:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return connector


def add_card(
    slide,
    number,
    x,
    y,
    w,
    h,
    title,
    formula,
    detail,
    accent,
    tag=None,
    formula_size=11.4,
    detail_size=8.9,
):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "panel", "line", 0.9)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, h, accent)
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.16, y + 0.15, 0.31, 0.31, accent)
    add_text(
        slide,
        str(number),
        x + 0.16,
        y + 0.15,
        0.31,
        0.31,
        8.7,
        "bg",
        True,
        FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        title,
        x + 0.57,
        y + 0.13,
        w - 0.70,
        0.34,
        12.0,
        "white",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        formula,
        x + 0.18,
        y + 0.55,
        w - 0.36,
        0.50,
        formula_size,
        accent,
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        x + 0.18,
        y + h - 0.40,
        w - 0.36,
        0.25,
        detail_size,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    if tag:
        add_pill(slide, tag, x + w - 0.79, y + h - 0.31, 0.63, 0.20, accent)

def add_symbol_legend(slide):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.60,
        5.55,
        6.95,
        0.85,
        "panel",
        "coral",
        0.9,
    )
    add_text(
        slide,
        "FORMULA / CODE",
        0.78,
        5.63,
        1.62,
        0.20,
        8.6,
        "coral",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "α_{t,e}：token t → expert e 的混合权重；乘在 activated 上",
        2.48,
        5.63,
        4.76,
        0.20,
        7.9,
        "white",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Gate: r=√softplus(xW_g^T); I=TopK(r+b); α=s·r[I]/Σ_{j∈I}r_j",
        0.78,
        5.90,
        6.48,
        0.18,
        7.4,
        "muted",
        False,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Code: for e∈E_r: idx,top=where(I=e); y.index_add_(idx, Expert(x[idx], α[idx,top]))",
        0.78,
        6.13,
        6.48,
        0.18,
        7.1,
        "muted",
        False,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def add_lane_label(slide, number, title, detail, y, accent):
    add_text(
        slide,
        number,
        0.60,
        y,
        0.34,
        0.22,
        9.0,
        accent,
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        title,
        0.98,
        y,
        2.35,
        0.22,
        9.3,
        "white",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        3.28,
        y,
        4.8,
        0.22,
        8.7,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.45, y + 0.10, 4.28, 0.015, "line")


def add_strategy_band(slide):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.60,
        3.15,
        12.13,
        0.27,
        "panel_2",
        "line",
        0.8,
    )
    add_text(
        slide,
        "EXPERT OWNERSHIP / LOAD",
        0.78,
        3.18,
        1.92,
        0.20,
        7.4,
        "blue",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "固定分片：core r 持有 E_r",
        2.86,
        3.18,
        2.58,
        0.20,
        8.0,
        "white",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "内存受限：U = |∪ I_t|；按 token 负载 / 内存分配活跃专家",
        5.62,
        3.18,
        6.76,
        0.20,
        7.9,
        "amber",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def add_mono_flow_box(
    slide,
    x,
    y,
    w,
    h,
    title,
    formula="",
    detail="",
    title_size=10.2,
    formula_size=8.6,
    detail_size=7.5,
):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel", "black", 1.0)
    add_text(
        slide,
        title,
        x + 0.16,
        y + 0.06,
        w - 0.32,
        0.18,
        title_size,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    formula_height = 0.0
    if formula:
        formula_height = 0.17 * (formula.count("\n") + 1)
        add_text(
            slide,
            formula,
            x + 0.16,
            y + 0.25,
            w - 0.32,
            formula_height,
            formula_size,
            "black",
            True,
            FONT,
            PP_ALIGN.CENTER,
            MSO_ANCHOR.MIDDLE,
        )
    if detail:
        detail_top = y + 0.25 + formula_height + 0.03
        detail_height = h - (detail_top - y) - 0.04
        add_text(
            slide,
            detail,
            x + 0.16,
            detail_top,
            w - 0.32,
            max(0.10, detail_height),
            detail_size,
            "gray",
            False,
            FONT_CJK,
            PP_ALIGN.CENTER,
            MSO_ANCHOR.MIDDLE,
        )


def add_horizontal_flow_box(
    slide,
    x,
    y,
    w,
    h,
    title,
    formula="",
    detail="",
    title_w=2.55,
    formula_w=3.95,
    title_size=10.0,
    formula_size=8.1,
    detail_size=7.2,
):
    """Draw a wide flow step with title, formula, and note on one row."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel", "black", 1.0)
    inner_x = x + 0.18
    inner_y = y + 0.07
    inner_h = h - 0.14
    available_w = w - 0.36
    actual_title_w = min(title_w, available_w * 0.34)
    actual_formula_w = min(formula_w, available_w - actual_title_w - 0.90)
    actual_detail_w = max(0.80, available_w - actual_title_w - actual_formula_w)
    add_text(
        slide,
        title,
        inner_x,
        inner_y,
        actual_title_w,
        inner_h,
        title_size,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    if formula:
        add_text(
            slide,
            formula,
            inner_x + actual_title_w + 0.08,
            inner_y,
            actual_formula_w,
            inner_h,
            formula_size,
            "black",
            True,
            FONT,
            PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )
    if detail:
        add_text(
            slide,
            detail,
            inner_x + actual_title_w + actual_formula_w + 0.16,
            inner_y,
            actual_detail_w - 0.16,
            inner_h,
            detail_size,
            "gray",
            False,
            FONT_CJK,
            PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )


def add_main_process_box(slide, x, y, w, h, text, size=16.0):
    """Draw one main-flow box whose only content is the process label."""
    shape = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel", "black", 1.1)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = Pt(10)
    frame.margin_right = Pt(10)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    for run in paragraph.runs:
        run.font.name = FONT_CJK
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = rgb(COLORS["black"])
    return shape


def add_external_note(
    slide,
    x,
    y,
    w,
    h,
    formula,
    detail="",
    formula_size=8.3,
    detail_size=7.5,
    align=PP_ALIGN.LEFT,
):
    """Add formula and explanation outside a main-flow box in one text frame."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.fill.background()
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = formula.splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = CODE_FONT
            run.font.size = Pt(formula_size)
            run.font.bold = True
            run.font.color.rgb = rgb(COLORS["black"])
    if detail:
        paragraph = frame.add_paragraph()
        paragraph.text = detail
        paragraph.alignment = align
        paragraph.space_before = Pt(1)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = FONT_CJK
            run.font.size = Pt(detail_size)
            run.font.bold = False
            run.font.color.rgb = rgb(COLORS["gray"])
    return box


def add_text_shape_box(
    slide,
    x,
    y,
    w,
    h,
    title,
    lines,
    title_size=10.0,
    body_size=7.8,
    body_font=FONT_CJK,
):
    """Draw a side explanation panel with all text inside its single shape."""
    shape = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel", "black", 1.0)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(10)
    frame.margin_right = Pt(10)
    frame.margin_top = Pt(7)
    frame.margin_bottom = Pt(6)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(6)
    for run in paragraph.runs:
        run.font.name = FONT_CJK
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = rgb(COLORS["black"])
    for line in lines:
        paragraph = frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(2)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.name = body_font
            run.font.size = Pt(body_size)
            run.font.bold = False
            run.font.color.rgb = rgb(COLORS["black"])
    return shape


def add_mono_down_arrow(slide, x, y1, y2, dashed=False):
    add_connector(slide, x, y1, x, y2 - 0.10, "black", 1.1, dashed)
    add_arrow(slide, x - 0.08, y2 - 0.16, 0.16, 0.16, "black", "down")


def add_compact_down_arrow(slide, x, y1, y2, dashed=False):
    """Draw a short arrow that leaves room for an external note."""
    add_connector(slide, x, y1, x, y2 - 0.06, "black", 1.0, dashed)
    add_arrow(slide, x - 0.06, y2 - 0.10, 0.12, 0.10, "black", "down")


def build_legacy_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.055, "cyan")
    add_text(
        slide,
        "DEEPSEEK V4 FLASH  /  MIXTURE-OF-EXPERTS",
        0.60,
        0.22,
        6.5,
        0.20,
        8.5,
        "cyan",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "多核 MoE：Token 选专家与跨核聚合",
        0.60,
        0.46,
        7.8,
        0.43,
        25.0,
        "white",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Score-based MoE · Expert Parallel · 一个 token 的完整路径",
        0.62,
        0.95,
        7.7,
        0.22,
        9.5,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_pill(slide, "完整模型  256E / TOP-6", 9.06, 0.37, 1.78, 0.33, "blue")
    add_pill(slide, "debug  16E / TOP-6", 10.98, 0.37, 1.75, 0.33, "amber")
    add_text(
        slide,
        "route_scale = 1.5",
        10.98,
        0.82,
        1.75,
        0.18,
        8.2,
        "amber",
        True,
        FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )

    add_lane_label(
        slide,
        "01",
        "ROUTE SELECTION  /  选择哪些专家",
        "全局选择 Top-K，生成 (token_pos, expert_id, α)",
        1.32,
        "cyan",
    )
    add_lane_label(
        slide,
        "02",
        "EXPERT EXECUTION  /  计算与聚合",
        "按专家组队 → 批量 FFN → 归位 → routed 汇总 + shared",
        3.45,
        "green",
    )

    top_y, top_w, top_h = 1.60, 2.10, 1.52
    top_x = [0.60, 3.10, 5.60, 8.10, 10.60]
    for left in top_x[:-1]:
        add_arrow(slide, left + top_w + 0.09, top_y + 0.61, 0.22, 0.18, "line")

    add_card(
        slide,
        1,
        top_x[0],
        top_y,
        top_w,
        top_h,
        "Token 表示",
        "x_t  ∈  R^D",
        "FFN RMSNorm 后  ·  T = B × S",
        "cyan",
    )
    add_card(
        slide,
        2,
        top_x[1],
        top_y,
        top_w,
        top_h,
        "Gate 线性投影",
        "z_t = x_t W_g^T",
        "每个 token 得到 E 个分数",
        "blue",
    )
    add_card(
        slide,
        3,
        top_x[2],
        top_y,
        top_w,
        top_h,
        "非线性打分",
        "r_t = √softplus(z_t)",
        "原始路由分数  ·  FP32",
        "green",
    )
    add_card(
        slide,
        4,
        top_x[3],
        top_y,
        top_w,
        top_h,
        "偏置 + Top-K",
        "q = r + b\nI = TopK(q, K=6)",
        "bias 只影响“选谁”",
        "amber",
        "SELECT",
        10.8,
    )
    add_card(
        slide,
        5,
        top_x[4],
        top_y,
        top_w,
        top_h,
        "路由权重",
        "α_{t,e} = s · r_{t,e}\n/ Σ_{j∈I_t} r_{t,j}",
        "s=1.5  ·  使用原始 r，不含 b",
        "coral",
        "WEIGHT",
        9.0,
    )

    add_strategy_band(slide)
    add_connector(slide, 11.65, 3.12, 11.65, 3.15, "cyan", 1.4)
    add_connector(slide, 11.65, 3.42, 11.65, 3.71, "cyan", 1.4)
    add_connector(slide, 11.65, 3.71, 1.65, 3.71, "cyan", 1.4)
    add_connector(slide, 1.65, 3.71, 1.65, 3.78, "cyan", 1.4)
    add_arrow(slide, 1.56, 3.69, 0.18, 0.18, "cyan", "down")

    bottom_y, bottom_w, bottom_h = 3.78, 2.10, 1.43
    bottom_x = [0.60, 3.10, 5.60, 8.10, 10.60]
    for left in bottom_x[:-1]:
        add_arrow(slide, left + bottom_w + 0.09, bottom_y + 0.59, 0.22, 0.18, "line")

    add_card(
        slide,
        6,
        bottom_x[0],
        bottom_y,
        bottom_w,
        bottom_h,
        "Dispatch / 组队",
        "(pos, x_t, e, α)",
        "按 expert_id 分桶；广播 / all-to-all",
        "cyan",
        formula_size=9.7,
        detail_size=7.7,
    )
    add_card(
        slide,
        7,
        bottom_x[1],
        bottom_y,
        bottom_w,
        bottom_h,
        "Expert batch",
        "W1/W3 → SwiGLU → α → W2",
        "批量执行；α × activated",
        "blue",
        formula_size=8.8,
        detail_size=7.6,
    )
    add_card(
        slide,
        8,
        bottom_x[2],
        bottom_y,
        bottom_w,
        bottom_h,
        "Local combine",
        "index_add(token_pos)",
        "按 token_pos 还原并累加",
        "green",
        formula_size=9.7,
        detail_size=8.0,
    )
    add_card(
        slide,
        9,
        bottom_x[3],
        bottom_y,
        bottom_w,
        bottom_h,
        "Routed AllReduce",
        "Σ_r y_r^routed",
        "routed 按 token 逐元素求和",
        "purple",
        "TP SUM",
        formula_size=9.8,
        detail_size=7.5,
    )
    add_card(
        slide,
        10,
        bottom_x[4],
        bottom_y,
        bottom_w,
        bottom_h,
        "Shared + output",
        "Σ_r y_r^routed + S(x)",
        "shared 执行；不参与归约",
        "amber",
        formula_size=8.6,
        detail_size=7.5,
    )

    add_connector(slide, 1.65, 3.12, 1.65, 3.28, "muted", 1.2, True)
    add_connector(slide, 1.65, 3.28, 2.90, 3.28, "muted", 1.2, True)
    add_connector(slide, 2.90, 3.28, 2.90, 5.43, "muted", 1.2, True)
    add_connector(slide, 2.90, 5.43, 8.95, 5.43, "muted", 1.2, True)
    add_connector(slide, 8.95, 5.43, 8.95, 5.63, "muted", 1.2, True)
    add_text(
        slide,
        "旁路：原始 token 直接进入 shared expert",
        2.00,
        5.26,
        3.08,
        0.20,
        8.3,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_symbol_legend(slide)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.90, 5.64, 2.10, 0.79, "panel_2", "muted", 0.8)
    add_text(
        slide,
        "Shared expert",
        8.06,
        5.75,
        1.78,
        0.22,
        10.0,
        "white",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "SwiGLU(x)  ·  always on",
        8.02,
        6.03,
        1.86,
        0.18,
        8.0,
        "muted",
        False,
        FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_connector(slide, 10.00, 6.03, 11.05, 5.22, "amber", 1.5)
    add_arrow(slide, 10.92, 5.20, 0.18, 0.18, "amber", "up")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.60, 6.47, 2.10, 0.42, "panel_2", "cyan", 0.8)
    add_text(
        slide,
        "MoE 输出 [T,D]  →  [B,S,D]  →  HC post",
        10.72,
        6.53,
        1.86,
        0.25,
        7.2,
        "cyan",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_connector(slide, 11.65, 5.21, 11.65, 6.43, "cyan", 1.4)
    add_arrow(slide, 11.56, 6.37, 0.18, 0.18, "cyan", "down")

    add_text(
        slide,
        "参考实现：routed experts 按 rank 固定分片；内存受限可按 U=unique(I) 动态换入。",
        0.60,
        6.82,
        9.65,
        0.25,
        8.4,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "source  inference/model.py  ·  export/export_rank0_prefill.py",
        0.60,
        7.18,
        8.0,
        0.16,
        6.9,
        "muted",
        False,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 11.55, 7.17, 1.18, 0.018, "cyan")

    prs.core_properties.title = "DeepSeek-V4 Flash 专家路由流程"
    prs.core_properties.subject = "Score-based MoE router flowchart"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide flowchart generated from the repository routing path."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def build_deck(output: Path) -> None:
    """Create a spacious reference-style flowchart with one box per main step."""
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_text(
        slide,
        "WD3：多核 MoE 激活计算（Token 选专家后）",
        0.42,
        0.12,
        12.45,
        0.40,
        22.0,
        "red",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "前提：X=[x_t]、I=[e_t,k]、A=[α_t,k] 已给定；流程从 dispatch 开始",
        0.45,
        0.56,
        12.25,
        0.18,
        8.8,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    main_x, main_w = 0.78, 8.90
    center_x = main_x + main_w / 2
    side_x, side_w = 10.02, 3.02

    add_text_shape_box(
        slide,
        0.38,
        0.84,
        2.10,
        0.60,
        "内存足够策略",
        ["core r 固定持有 E_r", "W_e 常驻"],
        title_size=9.4,
        body_size=7.5,
    )
    add_main_process_box(
        slide,
        2.76,
        0.84,
        6.42,
        0.60,
        "输入（已知）：Tokens + Top-K expert_id + α",
        size=14.5,
    )
    add_text_shape_box(
        slide,
        10.84,
        0.84,
        2.10,
        0.60,
        "内存不足策略",
        ["U=unique({e_t,k})", "按批动态加载"],
        title_size=9.4,
        body_size=7.5,
    )

    flow = [
        (1.66, 0.40, "准备本地专家集合", "E_r；W_e={W1^e,W2^e,W3^e}", "每个 core 只执行 E_r"),
        (2.38, 0.40, "Token 与路由表分发", "R={(t,e_t,k,α_t,k)}", "广播到各 core，也可 all-to-all"),
        (3.10, 0.40, "按 expert_id 组队", "G_{r,e}={(t,x_t,α_t,k):e_t,k=e}", "同一 expert_id 的 token 成 batch"),
        (3.82, 0.48, "本地 Expert 批量计算", "h=SiLU(xW1)⊙(xW3);  y=W2(α·h)", "α 乘在 activated 上"),
        (4.76, 0.40, "按 token_pos 归位", "y_r[t] += y_t,e", "index_add_(0, token_index, routed)"),
        (5.48, 0.40, "Routed 结果 AllReduce", "y_routed[t]=Σ_r y_r[t]", "逐 token、逐维求和；只汇总 routed"),
        (6.20, 0.42, "加 Shared Expert，输出结果", "y[t]=y_routed[t]+y_shared[t]", "[T,D] → [B,S,D] → HC post"),
    ]
    for y, h, title, formula, detail in flow:
        add_main_process_box(slide, main_x, y, main_w, h, title, size=15.2)
        note_h = 0.22 if y == 3.82 else 0.16
        add_external_note(
            slide,
            main_x + 0.16,
            y + h + 0.04,
            main_w - 0.32,
            note_h,
            formula,
            detail,
            formula_size=7.7 if y == 3.81 else 7.9,
            detail_size=7.0,
            align=PP_ALIGN.CENTER,
        )

    add_text_shape_box(
        slide,
        side_x,
        1.62,
        side_w,
        1.05,
        "实现代码",
        ["idx,top=where(I==e)", "routed=Expert(X[idx], A[idx,top])", "每个 expert 形成一个 batch"],
        title_size=10.5,
        body_size=7.0,
        body_font=CODE_FONT,
    )
    add_text_shape_box(
        slide,
        side_x,
        2.93,
        side_w,
        1.48,
        "组队示意",
        ["e_0 :  t_0   t_3", "e_1 :  t_1   t_2   t_4", "e_2 :  t_0   t_2", "同一 expert_id 一队"],
        title_size=10.5,
        body_size=7.2,
        body_font=CODE_FONT,
    )
    add_main_process_box(
        slide,
        side_x,
        5.28,
        side_w,
        0.50,
        "Shared Expert（always on）",
        size=10.8,
    )
    add_external_note(
        slide,
        side_x + 0.08,
        5.83,
        side_w - 0.16,
        0.28,
        "y_shared[t]=S(x_t)",
        "原始 token 旁路；不参与 routed 归约",
        formula_size=7.8,
        detail_size=6.9,
        align=PP_ALIGN.CENTER,
    )

    add_compact_down_arrow(slide, center_x, 1.44, 1.66)
    add_compact_down_arrow(slide, center_x, 2.22, 2.38)
    add_compact_down_arrow(slide, center_x, 2.94, 3.10)
    add_compact_down_arrow(slide, center_x, 3.66, 3.82)
    add_compact_down_arrow(slide, center_x, 4.52, 4.76)
    add_compact_down_arrow(slide, center_x, 5.38, 5.48)
    add_compact_down_arrow(slide, center_x, 6.10, 6.20)

    add_connector(slide, 1.43, 1.44, 2.35, 1.62, "black", 1.0)
    add_connector(slide, 11.89, 1.44, 8.30, 1.62, "black", 1.0)
    add_connector(slide, 9.68, 3.28, side_x, 3.28, "black", 1.0, True)
    add_connector(slide, side_x, 5.53, 9.68, 6.31, "black", 1.0, True)
    add_arrow(slide, 9.60, 6.23, 0.18, 0.18, "black", "left")

    add_text(
        slide,
        "主流程：每一步只保留一个方框；公式与实现说明放在框外，需要时再查看",
        0.50,
        7.16,
        9.60,
        0.17,
        7.6,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "inference/model.py · ScoreBasedMoE.forward",
        10.02,
        7.16,
        2.92,
        0.17,
        6.3,
        "gray",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "WD3 多核 MoE 激活计算流程"
    prs.core_properties.subject = "Top-K expert dispatch, expert grouping, local execution and AllReduce"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide flowchart with one box per main step and formulas outside the boxes."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def build_vertical_deck(output: Path) -> None:
    """Legacy vertical layout retained for comparison."""
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.035, "black")
    add_text(
        slide,
        "WD3：Token 选专家后的多核执行（组队可选）",
        0.55,
        0.20,
        12.20,
        0.38,
        22.0,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "前提：每个 token 的 Top-K expert_id 与混合权重 α 已给定；流程从 dispatch 开始",
        0.57,
        0.64,
        12.0,
        0.20,
        9.2,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    main_x, main_w = 3.00, 7.30
    center_x = main_x + main_w / 2

    add_mono_flow_box(
        slide,
        0.55,
        0.94,
        2.25,
        0.62,
        "内存足够",
        "core r 固定持有 E_r",
        "W_e 常驻",
        title_size=9.5,
        formula_size=8.2,
        detail_size=7.2,
    )
    add_mono_flow_box(
        slide,
        main_x,
        0.94,
        main_w,
        0.62,
        "输入（已知）",
        "X=[x_t]  ·  I=[e_t,k]  ·  A=[α_t,k]",
        "T tokens；I 与 A 已给定",
        title_size=10.0,
        formula_size=8.8,
        detail_size=7.5,
    )
    add_mono_flow_box(
        slide,
        10.55,
        0.94,
        2.25,
        0.62,
        "内存不足",
        "U=unique({e_t,k})",
        "按批动态加载",
        title_size=9.5,
        formula_size=8.2,
        detail_size=7.2,
    )

    add_mono_flow_box(
        slide,
        main_x,
        1.75,
        main_w,
        0.58,
        "每核确定本地专家集合，并准备参数",
        "E = ⋃_r E_r  ·  W_e={W1^e,W2^e,W3^e}",
        "每核只执行 E_r",
        title_size=9.5,
        formula_size=7.6,
        detail_size=7.2,
    )
    add_mono_flow_box(
        slide,
        main_x,
        2.43,
        main_w,
        0.58,
        "Tokens + 全局路由表，广播给所有 core",
        "R = {(t, e_t,k, α_t,k)}",
        "也可 all-to-all：只发给目标 core",
        title_size=9.7,
        formula_size=8.5,
        detail_size=7.2,
    )
    add_mono_flow_box(
        slide,
        10.55,
        2.05,
        2.25,
        0.90,
        "实现代码",
        "idx,top=where(I==e)\nrouted=Expert(X[idx],A[idx,top])",
        "每个 expert 形成一个 token batch",
        title_size=9.3,
        formula_size=6.8,
        detail_size=6.8,
    )
    add_mono_flow_box(
        slide,
        main_x,
        3.11,
        main_w,
        0.65,
        "按 expert_id 筛选、分桶，组成 Token 小队",
        "G_{r,e}={(t,x_t,α_t,k): e_t,k=e}",
        "按 expert_id 分桶；同一专家一队",
        title_size=9.7,
        formula_size=7.6,
        detail_size=7.2,
    )
    add_mono_flow_box(
        slide,
        main_x,
        3.91,
        main_w,
        0.84,
        "Token 小队 + 本地 Expert 批量计算",
        "h_t,e=SiLU(x_t W1^e) ⊙ (x_t W3^e)\ny_t,e=W2^e(α_t,e · h_t,e)",
        "α 乘在 activated 上",
        title_size=9.7,
        formula_size=7.5,
        detail_size=7.2,
    )
    add_mono_flow_box(
        slide,
        main_x,
        4.89,
        main_w,
        0.58,
        "Token 小队归位：按 token_pos 累加",
        "y_r[t] += y_t,e",
        "index_add_(0, token_index, routed)",
        title_size=9.7,
        formula_size=8.3,
        detail_size=7.2,
    )
    add_mono_flow_box(
        slide,
        main_x,
        5.63,
        main_w,
        0.58,
        "所有 core 对 routed 结果做 AllReduce(sum)",
        "y_routed[t] = Σ_r y_r[t]",
        "逐 token、逐维求和；只汇总 routed",
        title_size=9.7,
        formula_size=8.3,
        detail_size=7.2,
    )
    add_mono_flow_box(
        slide,
        main_x,
        6.37,
        main_w,
        0.62,
        "最终输出",
        "y[t] = y_routed[t] + y_shared[t]",
        "[T,D] → [B,S,D] → HC post",
        title_size=10.0,
        formula_size=8.5,
        detail_size=7.2,
    )

    add_mono_flow_box(
        slide,
        10.55,
        4.88,
        2.25,
        0.78,
        "Shared Expert",
        "y_shared[t] = S(x_t)",
        "原始 token 旁路；always on",
        title_size=9.3,
        formula_size=8.0,
        detail_size=7.0,
    )

    add_mono_down_arrow(slide, center_x, 1.56, 1.75)
    add_mono_down_arrow(slide, center_x, 2.33, 2.43)
    add_mono_down_arrow(slide, center_x, 3.01, 3.11)
    add_mono_down_arrow(slide, center_x, 3.76, 3.91)
    add_mono_down_arrow(slide, center_x, 4.75, 4.89)
    add_mono_down_arrow(slide, center_x, 5.47, 5.63)
    add_mono_down_arrow(slide, center_x, 6.21, 6.37)

    add_connector(slide, 1.68, 1.56, 4.05, 1.75, "black", 1.0)
    add_connector(slide, 11.67, 1.56, 9.25, 1.75, "black", 1.0)

    add_connector(slide, 10.30, 3.43, 11.67, 4.88, "black", 1.0, True)
    add_connector(slide, 11.67, 5.66, 11.67, 6.10, "black", 1.0, True)
    add_connector(slide, 11.67, 6.10, 10.30, 6.68, "black", 1.0, True)

    prs.core_properties.title = "WD3 Token 选专家后的多核执行流程"
    prs.core_properties.subject = "Top-K expert dispatch, grouping, local execution and AllReduce"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable monochrome top-down flowchart based on the reference image."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def build_horizontal_deck(output: Path) -> None:
    """Legacy horizontal layout retained for comparison."""
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_text(
        slide,
        "WD3：Token 选专家后的多核执行流程（组队可选）",
        0.45,
        0.13,
        12.40,
        0.40,
        20.5,
        "red",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "前提：每个 token 的 Top-K expert_id 与混合权重 α 已给定；流程从 dispatch 开始",
        0.48,
        0.57,
        12.25,
        0.18,
        8.7,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    main_x, main_w = 1.00, 9.20
    center_x = main_x + main_w / 2

    add_mono_flow_box(
        slide,
        0.45,
        0.86,
        2.35,
        0.64,
        "内存足够策略",
        "core r 固定持有 E_r",
        "W_e 常驻",
        title_size=9.3,
        formula_size=7.6,
        detail_size=6.9,
    )
    add_horizontal_flow_box(
        slide,
        3.00,
        0.86,
        7.20,
        0.64,
        "输入（已知）",
        "X=[x_t] · I=[e_t,k] · A=[α_t,k]",
        "T tokens；I、A 已给定",
        title_w=1.35,
        formula_w=3.95,
        title_size=10.0,
        formula_size=8.5,
        detail_size=7.4,
    )
    add_mono_flow_box(
        slide,
        10.55,
        0.86,
        2.35,
        0.64,
        "内存不足策略",
        "U=unique({e_t,k})",
        "按批动态加载",
        title_size=9.3,
        formula_size=7.6,
        detail_size=6.9,
    )

    flow = [
        (1.72, 0.50, "本地专家集合", "E_r；W_e={W1^e,W2^e,W3^e}", "每个 core 只执行 E_r"),
        (2.44, 0.50, "Token + 路由表", "R={(t,e_t,k,α_t,k)}", "广播到各 core / all-to-all"),
        (3.16, 0.56, "按 expert_id 组队", "G_{r,e}={(t,x_t,α_t,k):e_t,k=e}", "同一专家的 token 成 batch"),
        (3.94, 0.82, "专家批量计算", "h=SiLU(xW1)⊙(xW3)\ny=W2(α·h)", "α × activated；Expert(X[idx], A[idx,top])"),
        (4.98, 0.50, "Token 归位", "y_r[t] += y_t,e", "按 token_pos 累加；index_add_(...)"),
        (5.70, 0.50, "Routed AllReduce", "y_routed[t]=Σ_r y_r[t]", "逐 token、逐维求和"),
        (6.42, 0.55, "最终输出", "y[t]=y_routed[t]+y_shared[t]", "[T,D] → [B,S,D] → HC post"),
    ]
    for y, h, title, formula, detail in flow:
        add_horizontal_flow_box(
            slide,
            main_x,
            y,
            main_w,
            h,
            title,
            formula,
            detail,
            title_w=2.35,
            formula_w=4.05,
            title_size=10.0,
            formula_size=7.6 if "\n" in formula else 8.3,
            detail_size=7.1,
        )

    add_shape(slide, MSO_SHAPE.RECTANGLE, 10.55, 1.76, 2.35, 1.16, "panel", "black", 1.0)
    add_text(
        slide,
        "实现代码",
        10.72,
        1.84,
        2.01,
        0.18,
        9.4,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "idx,top=where(I==e)\nrouted=Expert(X[idx], A[idx,top])",
        10.72,
        2.10,
        2.01,
        0.36,
        6.5,
        "black",
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "每个 expert 形成一个 batch",
        10.72,
        2.62,
        2.01,
        0.17,
        6.8,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_shape(slide, MSO_SHAPE.RECTANGLE, 10.55, 3.10, 2.35, 1.35, "panel", "black", 1.0)
    add_text(
        slide,
        "组队示意",
        10.72,
        3.18,
        2.01,
        0.18,
        9.4,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "e_0 :  t_0   t_3\ne_1 :  t_1   t_2   t_4\ne_2 :  t_0   t_2",
        10.72,
        3.48,
        2.01,
        0.54,
        7.2,
        "black",
        False,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "同一 expert_id 一队",
        10.72,
        4.14,
        2.01,
        0.17,
        6.8,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_mono_flow_box(
        slide,
        10.55,
        5.18,
        2.35,
        0.72,
        "Shared Expert",
        "y_shared[t]=S(x_t)",
        "always on；不参与归约",
        title_size=9.2,
        formula_size=7.8,
        detail_size=6.8,
    )

    add_mono_down_arrow(slide, center_x, 1.50, 1.72)
    add_mono_down_arrow(slide, center_x, 2.22, 2.44)
    add_mono_down_arrow(slide, center_x, 2.94, 3.16)
    add_mono_down_arrow(slide, center_x, 3.72, 3.94)
    add_mono_down_arrow(slide, center_x, 4.76, 4.98)
    add_mono_down_arrow(slide, center_x, 5.48, 5.70)
    add_mono_down_arrow(slide, center_x, 6.20, 6.42)

    add_connector(slide, 1.63, 1.50, 2.60, 1.72, "black", 1.0)
    add_connector(slide, 11.72, 1.50, 8.90, 1.72, "black", 1.0)
    add_mono_down_arrow(slide, 10.38, 1.50, 5.18, dashed=True)
    add_connector(slide, 10.38, 5.90, 10.20, 6.68, "black", 1.0, True)

    add_text(
        slide,
        "输入 I、A 已知后，主流程只处理 dispatch、专家计算、归位与跨核聚合",
        0.55,
        7.11,
        9.90,
        0.17,
        7.7,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "inference/model.py  ·  ScoreBasedMoE.forward",
        10.55,
        7.11,
        2.35,
        0.17,
        6.4,
        "gray",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "WD3 Token 选专家后的多核执行流程"
    prs.core_properties.subject = "Horizontal monochrome MoE dispatch and aggregation flowchart"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide flowchart with horizontal text layout based on the reference image."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


if __name__ == "__main__":
    target = Path(__file__).with_name("moe_route_flowchart.pptx")
    build_deck(target)
    print(target)
