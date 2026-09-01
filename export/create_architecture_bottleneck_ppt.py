#!/usr/bin/env python3
"""Create a one-slide Chinese translation of the architecture co-design diagram."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from create_moe_route_ppt import (
    COLORS,
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_connector,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5


def add_dashed_frame(slide, x, y, w, h):
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame.fill.background()
    frame.line.color.rgb = rgb(COLORS["line"])
    frame.line.width = Pt(1.0)
    frame.line.dash_style = 1
    return frame


def add_bullet(slide, x, y, title, accent, width=2.45):
    add_shape(slide, MSO_SHAPE.OVAL, x, y + 0.05, 0.12, 0.12, accent)
    add_text(
        slide,
        title,
        x + 0.22,
        y,
        width - 0.22,
        0.25,
        12.0,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def add_small_node(slide, x, y, w, h, text, fill, line, size=11.0):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line, 1.0)
    add_text(
        slide,
        text,
        x + 0.08,
        y + 0.05,
        w - 0.16,
        h - 0.10,
        size,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_architecture_column(slide, x, label, index):
    add_small_node(slide, x, 1.78, 1.52, 0.46, f"HBM 切片 {index}", "panel_2", "blue", 10.4)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.70, 2.24, 0.12, 0.22, "green")
    add_small_node(slide, x + 0.27, 2.46, 0.98, 0.42, "本地视图", "panel", "green", 9.4)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.70, 2.88, 0.12, 0.28, "green")
    add_small_node(slide, x, 3.16, 1.52, 0.58, f"计算核心切片 {index}", "light", "coral", 10.4)
    add_connector(slide, x + 0.76, 3.74, x + 0.76, 4.08, "coral", 1.4)
    add_connector(slide, x + 0.76, 4.08, x + 0.76, 4.36, "line", 1.4)


def add_section(slide, x, y, title, body, accent, body_width=2.70):
    add_text(
        slide,
        title,
        x,
        y,
        body_width,
        0.28,
        16.0,
        accent,
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        body,
        x,
        y + 0.34,
        body_width,
        0.68,
        11.0,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_text(
        slide,
        "协同设计架构，解决系统瓶颈",
        0.33,
        0.16,
        9.30,
        0.42,
        25.0,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.28, 0.10, 3.42, 0.54, "light")
    add_text(
        slide,
        "• Gather：维度变大\n• Reduce：维度不变",
        8.48,
        0.14,
        3.02,
        0.46,
        10.6,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_dashed_frame(slide, 0.28, 1.34, 9.15, 4.57)
    add_bullet(slide, 0.48, 0.83, "张量 Gather（片上）", "black", 2.20)
    add_bullet(slide, 5.75, 0.83, "AllGather（片间）", "black", 2.20)
    add_bullet(slide, 10.86, 0.83, "量化（FM 在线）", "black", 2.08)

    add_architecture_column(slide, 0.72, "HBM", 1)
    add_architecture_column(slide, 2.90, "HBM", 2)
    add_architecture_column(slide, 5.08, "HBM", 3)
    add_architecture_column(slide, 7.26, "HBM", 64)
    add_text(slide, "…", 6.42, 2.82, 0.46, 0.38, 25.0, "gray", True, FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 4.36, 8.28, 0.48, "panel_2", "amber", 1.2)
    add_text(
        slide,
        "集合通信网络 · 高带宽、低时延",
        0.98,
        4.47,
        7.74,
        0.24,
        12.0,
        "gray",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 4.88, 8.28, 0.38, "light", "line", 1.0)
    add_text(
        slide,
        "通用 NoC · 带宽较低、时延较高 · 更灵活",
        0.91,
        4.96,
        7.90,
        0.20,
        10.6,
        "gray",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_connector(slide, 9.00, 5.07, 9.74, 5.07, "amber", 2.0)
    add_small_node(slide, 9.73, 4.78, 1.52, 0.74, "Scale-up\n以太网桥", "light", "amber", 10.2)

    add_section(
        slide,
        9.92,
        1.50,
        "最快本地路径",
        "将每个计算核心切片与一个 HBM 切片配对，形成低时延、\n高带宽的本地访问视图。",
        "green",
        3.03,
    )
    add_section(
        slide,
        9.92,
        2.78,
        "专用核间互连",
        "集合通信网络针对常见通信模式优化带宽和时延。",
        "amber",
        3.03,
    )
    add_section(
        slide,
        9.92,
        3.72,
        "保持灵活性",
        "共享 NoC 支持通用通信，并连接 Scale-up 网络。",
        "black",
        3.03,
    )

    add_text(
        slide,
        "将常用操作数保留在本地；仅在必要时使用共享路径",
        0.76,
        6.45,
        11.88,
        0.38,
        21.0,
        "green",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "显式布局和分布式控制，避免数据搬移与同步主导执行",
        0.62,
        7.02,
        12.10,
        0.24,
        13.0,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "协同设计架构，解决系统瓶颈"
    prs.core_properties.subject = "Chinese translation of HBM, core slice and collective network architecture diagram"
    prs.core_properties.author = "DeepSeek V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide Chinese redraw based on the supplied reference image."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("architecture_bottleneck_cn.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()