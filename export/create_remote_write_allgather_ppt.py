#!/usr/bin/env python3
"""Create a one-slide eight-core remote-write AllGather flowchart."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from create_moe_route_ppt import (
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5


def add_step_box(
    slide,
    text,
    x,
    y,
    w,
    h,
    accent,
    font=FONT_CJK,
    size=8.3,
    bold=True,
    fill="panel",
):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, accent, 0.75)
    add_text(
        slide,
        text,
        x + 0.06,
        y + 0.03,
        w - 0.12,
        h - 0.06,
        size,
        "black",
        bold,
        font,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_arrow_between(slide, x, y, h, size=11.0):
    add_text(
        slide,
        "→",
        x,
        y,
        0.05,
        h,
        size,
        "muted",
        True,
        FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.045, "cyan")
    add_text(
        slide,
        "8 核远程写入版 AllGather：算子执行流程",
        0.40,
        0.16,
        9.60,
        0.43,
        21.5,
        "red",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "八个核心同时启动同一个 collective operator；每个源分片写入其他 7 个核心的远程接收槽位",
        0.43,
        0.63,
        12.25,
        0.20,
        9.1,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.36, 0.96, 12.61, 0.34, "panel_2", "line", 0.8)
    add_text(
        slide,
        "核心语义",
        0.54,
        1.03,
        0.78,
        0.19,
        8.4,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Y_j[i] = X_i",
        1.44,
        1.03,
        1.14,
        0.19,
        9.0,
        "cyan",
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "目标核 j 的第 i 个接收槽位",
        2.67,
        1.03,
        2.32,
        0.19,
        8.2,
        "white",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "AllGather = 搬运 + 拼接，不做数值相加",
        5.19,
        1.03,
        3.40,
        0.19,
        8.3,
        "blue",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "C = S / 8",
        11.31,
        1.03,
        1.38,
        0.19,
        8.4,
        "amber",
        True,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    table_x = 0.36
    table_y = 1.39
    table_w = 12.61
    table_h = 4.76
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, table_x, table_y, table_w, table_h, "panel", "line", 0.8)

    # Column geometry is shared by the header and all eight core rows.
    column_widths = [0.66, 0.90, 0.90, 0.88, 3.00, 1.17, 0.95, 3.64]
    gap = 0.05
    column_x = []
    cursor = table_x + 0.08
    for width in column_widths:
        column_x.append(cursor)
        cursor += width + gap

    headers = [
        "CORE",
        "LAUNCH",
        "LOCAL",
        "OWN SLOT",
        "REMOTE WRITE / NoC",
        "ORDER + FLAG",
        "WAIT",
        "LOCAL OUTPUT",
    ]
    header_y = table_y + 0.07
    header_h = 0.24
    for index, header in enumerate(headers):
        add_text(
            slide,
            header,
            column_x[index],
            header_y,
            column_widths[index],
            header_h,
            7.0 if index not in (4, 7) else 6.7,
            "muted",
            True,
            FONT,
            PP_ALIGN.CENTER,
            MSO_ANCHOR.MIDDLE,
        )

    row_y = table_y + 0.43
    row_h = 0.43
    row_step = 0.53
    accents = ["cyan", "blue", "green", "amber", "coral", "purple", "cyan", "blue"]
    for core in range(8):
        y = row_y + core * row_step
        accent = accents[core]
        if core % 2 == 1:
            add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                table_x + 0.05,
                y - 0.04,
                table_w - 0.10,
                row_h + 0.08,
                "light",
            )

        add_step_box(
            slide,
            f"Core {core}",
            column_x[0],
            y,
            column_widths[0],
            row_h,
            accent,
            FONT,
            8.0,
        )
        add_step_box(
            slide,
            "start\nAllGather",
            column_x[1],
            y,
            column_widths[1],
            row_h,
            "blue",
            FONT_CJK,
            7.3,
        )
        add_step_box(
            slide,
            f"X{core}\n[C]",
            column_x[2],
            y,
            column_widths[2],
            row_h,
            "cyan",
            CODE_FONT,
            8.0,
        )
        add_step_box(
            slide,
            f"Y{core}[{core}]\n= X{core}",
            column_x[3],
            y,
            column_widths[3],
            row_h,
            "green",
            CODE_FONT,
            7.0,
        )
        add_step_box(
            slide,
            f"remote_write: X{core} → Yj[{core}]\nj ≠ {core}  ·  7 × put",
            column_x[4],
            y,
            column_widths[4],
            row_h,
            "cyan",
            CODE_FONT,
            7.1,
            fill="panel_2",
        )
        add_step_box(
            slide,
            f"fence → ready_j[{core}]",
            column_x[5],
            y,
            column_widths[5],
            row_h,
            "amber",
            CODE_FONT,
            7.0,
        )
        add_step_box(
            slide,
            f"wait\nready_{core}[0:8]",
            column_x[6],
            y,
            column_widths[6],
            row_h,
            "purple",
            CODE_FONT,
            6.9,
        )
        add_step_box(
            slide,
            f"Y{core} = [X0, X1, …, X7]",
            column_x[7],
            y,
            column_widths[7],
            row_h,
            "green",
            CODE_FONT,
            8.0,
            fill="panel_2",
        )

        arrow_y = y + 0.02
        for index in range(7):
            add_arrow_between(
                slide,
                column_x[index] + column_widths[index] + 0.005,
                arrow_y,
                row_h - 0.04,
                9.0,
            )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.36, 6.34, 7.28, 0.48, "panel_2", "green", 0.9)
    add_text(
        slide,
        "最终：Y0 = Y1 = … = Y7 = [X0, X1, …, X7]",
        0.58,
        6.47,
        6.84,
        0.20,
        10.0,
        "green",
        True,
        CODE_FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.83, 6.34, 5.14, 0.48, "panel_2", "amber", 0.9)
    add_text(
        slide,
        "每核：7×C 发送 + 7×C 接收；全系统 56×C = 7S",
        8.03,
        6.47,
        4.74,
        0.20,
        8.0,
        "amber",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "硬件最小能力：remote DMA / NoC + 远程地址偏移 + 写入顺序/fence + completion flag；无需 Reduce Engine",
        0.40,
        6.98,
        12.45,
        0.20,
        8.0,
        "muted",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "AllGather remote-write flow  ·  8 cores  ·  editable diagram",
        0.40,
        7.25,
        6.30,
        0.14,
        6.7,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "8 核远程写入版 AllGather 算子执行流程"
    prs.core_properties.subject = "Eight-core remote-memory-write AllGather flowchart"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide flowchart showing simultaneous operator execution on eight cores."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("remote_write_allgather_8core.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()