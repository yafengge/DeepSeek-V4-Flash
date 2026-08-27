#!/usr/bin/env python3
"""Create a one-slide recommended eight-core Broadcast flowchart."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from create_moe_route_ppt import (
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5


def add_process_box(
    slide,
    y,
    h,
    title,
    formula,
    title_width=2.78,
    title_size=10.7,
    formula_size=8.5,
):
    x, w = 0.48, 7.72
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, "panel", "black", 1.2)
    add_text(
        slide,
        title,
        x + 0.15,
        y + 0.05,
        title_width,
        h - 0.10,
        title_size,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        formula,
        x + title_width + 0.15,
        y + 0.05,
        w - title_width - 0.30,
        h - 0.10,
        formula_size,
        "black",
        False,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )


def add_endpoint(slide, y, h, title, detail):
    x, w = 1.55, 5.58
    add_shape(slide, MSO_SHAPE.OVAL, x, y, w, h, "panel", "black", 1.35)
    add_text(
        slide,
        title,
        x + 0.14,
        y + 0.04,
        w - 0.28,
        0.22,
        11.6,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        x + 0.16,
        y + 0.25,
        w - 0.32,
        h - 0.30,
        8.0,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_down_arrow(slide, y):
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, 4.10, y, 0.25, 0.14, "black")


def add_capability(slide, number, title, detail, y):
    add_text(
        slide,
        str(number),
        8.86,
        y,
        0.27,
        0.25,
        9.0,
        "black",
        True,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        title,
        9.23,
        y,
        3.16,
        0.25,
        8.7,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        detail,
        9.23,
        y + 0.23,
        3.16,
        0.31,
        7.5,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_text(
        slide,
        "8 核 Broadcast 推荐执行流程",
        0.48,
        0.18,
        10.4,
        0.40,
        20.0,
        "black",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "一个源 Core，多个目标 Core：NoC 负责一源多目的复制，目标端写入本地内存",
        0.51,
        0.61,
        11.80,
        0.20,
        8.9,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Y_i = X_0   (i=0,…,7)",
        10.35,
        0.61,
        2.25,
        0.20,
        8.5,
        "black",
        True,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    add_text(
        slide,
        "推荐数据流 / Core_i 视角",
        0.48,
        0.83,
        7.72,
        0.17,
        7.6,
        "gray",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "NoC 必备能力",
        8.66,
        0.83,
        4.04,
        0.17,
        7.6,
        "gray",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_endpoint(
        slide,
        1.02,
        0.34,
        "开始",
        "Core_0 … Core_7 同时启动 Broadcast(root=0)",
    )
    add_down_arrow(slide, 1.39)
    add_process_box(
        slide,
        1.51,
        0.34,
        "准备目标缓冲区",
        "all Core: B_i[S]；sync_prepare()",
        formula_size=8.4,
    )
    add_down_arrow(slide, 1.88)
    add_process_box(
        slide,
        2.00,
        0.39,
        "源 Core 提交一次 BCAST",
        "Core_0: src=X_0；dst={1,…,7}；bytes=S",
        title_width=2.96,
        formula_size=8.0,
    )
    add_down_arrow(slide, 2.43)
    add_process_box(
        slide,
        2.55,
        0.39,
        "NoC 读取并分包",
        "DMA read(X_0) → packetize；源端只注入 1 次",
        title_width=2.52,
        formula_size=7.8,
    )
    add_down_arrow(slide, 2.98)
    add_process_box(
        slide,
        3.10,
        0.42,
        "NoC 组播路由与复制",
        "multicast tree：在分支处复制 packet/flit",
        title_width=2.72,
        formula_size=8.0,
    )
    add_down_arrow(slide, 3.55)
    add_process_box(
        slide,
        3.67,
        0.42,
        "目标端点远程写入",
        "Core_j endpoint: B_j ← X_0；local_done_j=1",
        title_width=2.72,
        formula_size=7.9,
    )
    add_down_arrow(slide, 4.12)
    add_process_box(
        slide,
        4.24,
        0.39,
        "completion / 完成确认",
        "completion: ACK / done_bitmap = 11111110",
        title_width=1.65,
        formula_size=8.7,
    )
    add_down_arrow(slide, 4.68)
    add_process_box(
        slide,
        4.80,
        0.38,
        "顺序与集合同步",
        "fence core_0 → sync(core_0,…,core_7)",
        title_width=2.22,
        formula_size=8.0,
    )
    add_down_arrow(slide, 5.27)
    add_process_box(
        slide,
        5.39,
        0.35,
        "所有 Core 使用本地副本",
        "B_i = X_0",
        title_width=2.82,
        formula_size=9.3,
    )
    add_down_arrow(slide, 5.80)
    add_endpoint(
        slide,
        5.94,
        0.37,
        "结束",
        "Broadcast 完成",
    )

    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.66, 1.10, 4.04, 6.18, "light", "black", 1.0)
    add_text(
        slide,
        "数据面 + 控制面",
        8.86,
        1.34,
        3.50,
        0.22,
        10.3,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    capabilities = [
        (1, "Broadcast 事务头", "root / dst_group / src / dst_offset / bytes / seq"),
        (2, "组播路由与复制", "路由树；分支复制；避免重复包和环路"),
        (3, "remote DMA + endpoint", "源读、burst 传输、目标本地 SRAM/DRAM 写入"),
        (4, "FIFO 与流控", "credit / backpressure；处理慢目标和拥塞"),
        (5, "完成与错误反馈", "DMA done、ACK bitmap、timeout、error code"),
        (6, "顺序与内存可见性", "payload → fence → ready；release/acquire、coherence"),
    ]
    capability_y = 1.82
    for capability in capabilities:
        add_capability(slide, *capability, capability_y)
        capability_y += 0.76
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.86, 6.45, 3.50, 0.012, "black")
    add_text(
        slide,
        "无需 Reduce Engine：Broadcast 只复制数据，不做数值相加。",
        8.86,
        6.65,
        3.50,
        0.44,
        8.2,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        "软件提交 1 个 BCAST；NoC 内部完成复制与转发。",
        8.86,
        7.09,
        3.50,
        0.22,
        7.9,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "8 核 Broadcast 推荐执行流程与 NoC 能力"
    prs.core_properties.subject = "Recommended remote-memory broadcast flow and NoC requirements"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide black-and-white Broadcast flowchart with NoC capability annotations."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("broadcast_noc_recommended_flow.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()