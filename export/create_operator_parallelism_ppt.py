#!/usr/bin/env python3
"""Create a one-slide overview of the operator parallelism strategy."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from create_moe_route_ppt import (
    COLORS,
    CODE_FONT,
    FONT,
    FONT_CJK,
    add_pill,
    add_shape,
    add_text,
    rgb,
)


W, H = 13.333, 7.5


def add_parallel_card(
    slide,
    number,
    x,
    y,
    w,
    h,
    title,
    tag,
    conclusion,
    body,
    accent,
    tag_width=1.25,
):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "panel", "line", 0.9)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, h, accent)
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.17, y + 0.15, 0.31, 0.31, accent)
    add_text(
        slide,
        str(number),
        x + 0.17,
        y + 0.15,
        0.31,
        0.31,
        8.7,
        "bg",
        True,
        FONT,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        title,
        x + 0.59,
        y + 0.13,
        w - 0.78,
        0.34,
        11.4,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_pill(slide, tag, x + 0.19, y + 0.55, tag_width, 0.23, accent)
    add_text(
        slide,
        conclusion,
        x + 0.19,
        y + 0.84,
        w - 0.38,
        0.30,
        9.1,
        accent,
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        body,
        x + 0.19,
        y + 1.18,
        w - 0.38,
        h - 1.30,
        8.0,
        "white",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
    )


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.045, "cyan")
    add_text(
        slide,
        "六类算子的并行策略总览",
        0.48,
        0.22,
        8.6,
        0.43,
        22.0,
        "red",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "基于当前 TP=8 参考实现：先看参数/激活布局，再看每个 rank 的本地工作与跨 rank 通信",
        0.51,
        0.70,
        12.25,
        0.20,
        9.3,
        "gray",
        False,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.48, 1.00, 12.37, 0.40, "panel_2", "line", 0.8)
    add_text(
        slide,
        "读图：",
        0.67,
        1.09,
        0.56,
        0.20,
        8.8,
        "black",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Shard = 按维度/ID 分片",
        1.29,
        1.09,
        2.05,
        0.20,
        8.5,
        "blue",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Replicate = 每 rank 完整副本",
        3.50,
        1.09,
        2.65,
        0.20,
        8.5,
        "green",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "AllReduce = 求和",
        6.30,
        1.09,
        1.62,
        0.20,
        8.5,
        "purple",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "AllGather = 拼接",
        8.08,
        1.09,
        1.70,
        0.20,
        8.5,
        "coral",
        True,
        FONT_CJK,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "V= vocab_size  ·  D=hidden dim  ·  H=hc_mult",
        9.95,
        1.09,
        2.60,
        0.20,
        7.8,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    margin, gap = 0.48, 0.18
    card_w = (W - 2 * margin - 2 * gap) / 3
    card_x = [margin, margin + card_w + gap, margin + 2 * (card_w + gap)]
    card_h = 2.28
    top_y, bottom_y = 1.58, 4.08

    add_parallel_card(
        slide,
        1,
        card_x[0],
        top_y,
        card_w,
        card_h,
        "VocabParallelEmbedding",
        "VOCAB SHARD",
        "按词表行切分；局部 lookup 后求和",
        "参数：W_r=[V/TP,D]；input_ids 在各 rank 全量\nRank r：非本地 token mask=0，再查本地表\n通信：AllReduce(sum) → [B,S,D] 全量",
        "cyan",
        1.32,
    )
    add_parallel_card(
        slide,
        2,
        card_x[1],
        top_y,
        card_w,
        card_h,
        "Attention",
        "Q/O SHARD",
        "Q heads、O factors 切分；KV 每 rank 复制",
        "参数：wq_b/wo_a/wo_b 分片；wq_a、wkv 复制\nRank r：局部 Q heads 与 O groups；KV/cache 全量\n通信：输出 AllReduce；ratio-4 Indexer score 也归约",
        "blue",
        1.15,
    )
    add_parallel_card(
        slide,
        3,
        card_x[2],
        top_y,
        card_w,
        card_h,
        "Expert parallel MoE",
        "EXPERT-ID SHARD",
        "routed experts 按 expert_id 分片",
        "参数：E/TP 个 routed experts；gate/shared 复制\nRank r：完整 token + 路由表，只算本地 experts\n通信：routed AllReduce；shared 复制计算、不归约",
        "green",
        1.55,
    )
    add_parallel_card(
        slide,
        4,
        card_x[0],
        bottom_y,
        card_w,
        card_h,
        "RMSNorm",
        "REPLICATED",
        "完整 hidden 维度本地算；不是 hidden-shard TP",
        "参数/输入：γ 与 [B,S,D] 都是完整副本\nRank r：沿 D 本地计算 mean(x²)，再乘 γ\n通信：无；mean(-1) 不是跨 rank Reduce",
        "amber",
        1.05,
    )
    add_parallel_card(
        slide,
        5,
        card_x[1],
        bottom_y,
        card_w,
        card_h,
        "mHC (pre-mix + post-mix)",
        "HC REPLICATED",
        "沿 HC 维度混合；mHC 本身不切 TP",
        "参数：HC mixing/Sinkhorn 参数复制\nRank r：pre [B,S,H,D]→[B,S,D]；post 反向扩展\n通信：无；Attention/MoE 内部再执行 TP",
        "purple",
        1.53,
    )
    add_parallel_card(
        slide,
        6,
        card_x[2],
        bottom_y,
        card_w,
        card_h,
        "LM Head",
        "VOCAB SHARD",
        "按词表行切 logits；最后拼成完整 vocab",
        "参数：W_r=[V/TP,D]；hidden/Norm 在各 rank 全量\nRank r：last token → local logits [B,V/TP]\n通信：AllGather(concat) → [B,V]；权重独立于 Embedding",
        "coral",
        1.32,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.48, 6.56, 12.37, 0.43, "panel_2", "line", 0.8)
    add_text(
        slide,
        "COLLECTIVE PATH",
        0.67,
        6.66,
        1.42,
        0.20,
        8.2,
        "black",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Embedding  AllReduce  →  Attention  AllReduce  →  MoE  AllReduce  →  LM Head  AllGather",
        2.22,
        6.66,
        7.36,
        0.20,
        8.4,
        "blue",
        True,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "RMSNorm / mHC：local only",
        9.84,
        6.66,
        2.72,
        0.20,
        8.2,
        "green",
        True,
        FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "source  inference/model.py  ·  export/export_rank0_prefill.py",
        0.51,
        7.13,
        7.6,
        0.16,
        6.8,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "TP=world_size=8  ·  当前参考实现",
        9.80,
        7.13,
        3.05,
        0.16,
        6.8,
        "muted",
        False,
        CODE_FONT,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    prs.core_properties.title = "六类算子的并行策略总览"
    prs.core_properties.subject = "Tensor parallelism, expert parallelism and replicated operators"
    prs.core_properties.author = "DeepSeek-V4 Flash repository"
    prs.core_properties.comments = "Editable one-slide operator parallelism overview based on the reference implementation."
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    target = Path(__file__).with_name("operator_parallelism_overview.pptx")
    build_deck(target)
    print(target)


if __name__ == "__main__":
    main()